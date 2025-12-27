# ai-testpy_VLLM_TEACHER_ANS.py
# Stable server with:
# - legacy flat domains supported
# - nested Courses/<Course>/YEAR_X/SEM_Y/<Lesson> supported
# - year-scope escalation via scope="year"
# - ✅ accepts year alias: Courses__<Course>__YEAR_X
# - safe file serving for nested domain folders

from flask import Flask, request, jsonify, abort, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import safe_join
import os, io, re, time, json, hashlib, shutil, unicodedata
import numpy as np
import requests
import PyPDF2

try:
    import faiss
except Exception:
    faiss = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

try:
    import docx
except Exception:
    docx = None

try:
    import pptx
except Exception:
    pptx = None

try:
    import xlrd
except Exception:
    xlrd = None

try:
    import openpyxl
except Exception:
    openpyxl = None

app = Flask(__name__)
CORS(app)

# ============
# Config
# ============
DATA_ROOT = os.getenv("DATA_ROOT", "./domains")
CACHE_ROOT = os.getenv("CACHE_ROOT", "./.cache")
INDEX_CACHE_DIR = os.path.join(CACHE_ROOT, "indexes")

os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

API_KEY = os.getenv("API_KEY", "changeme")

CACHE_ADMIN_KEY = os.getenv("CACHE_ADMIN_KEY", "changeme-admin")

EMBEDDING_DEVICE = os.getenv("EMBEDDING_DEVICE", "cpu")

VLLM_URL = os.getenv("VLLM_URL", "http://127.0.0.1:8000/v1/chat/completions")
VLLM_MODEL = os.getenv("VLLM_MODEL", "meta-llama/Meta-Llama-3.1-8B-Instruct")
VLLM_TIMEOUT = int(os.getenv("VLLM_TIMEOUT", "120"))

MODEL_NAME = os.getenv("MODEL_NAME", "vLLM-chat")

TEMPERATURE = float(os.getenv("TEMPERATURE", "0.2"))
TOP_P = float(os.getenv("TOP_P", "0.9"))
MAX_TOKENS = int(os.getenv("MAX_TOKENS", "140"))

# vLLM token budgeting (prevents 400s when prompt is large)
VLLM_MAX_CONTEXT = int(os.getenv("VLLM_MAX_CONTEXT", "2048"))
VLLM_SAFETY_MARGIN = int(os.getenv("VLLM_SAFETY_MARGIN", "256"))
VLLM_MIN_COMPLETION = int(os.getenv("VLLM_MIN_COMPLETION", "96"))
VLLM_MAX_COMPLETION_CAP = int(os.getenv("VLLM_MAX_COMPLETION_CAP", str(MAX_TOKENS)))

# Retrieval tuning
CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "280"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "20"))
R_TOP_K = int(os.getenv("R_TOP_K", "4"))
NUM_ALTERNATE_QUERIES = int(os.getenv("NUM_ALTERNATE_QUERIES", "1"))
MIN_SIM_THRESHOLD = float(os.getenv("MIN_SIM_THRESHOLD", "0.30"))

TEACHER_MIN_SIM = float(os.getenv("TEACHER_MIN_SIM", "0.55"))

# Teacher answers cache:
# { domain_id: { "questions": [...], "answers": [...],
#               "embeddings": np.ndarray, "pages": [...], "file": "..." } }
teacher_answers_cache = {}

STRICT_MODE = os.getenv("STRICT_MODE", "true").lower() in ("1", "true", "yes")
EXTRACTIVE_ONLY = os.getenv("EXTRACTIVE_ONLY", "false").lower() in ("1", "true", "yes")

EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "64"))
MIN_WORDS_TO_INDEX = int(os.getenv("MIN_WORDS_TO_INDEX", "20"))

# Index cache
USE_INDEX_CACHE = os.getenv("USE_INDEX_CACHE", "true").lower() in ("1", "true", "yes")
INDEX_CACHE_VERSION = "v4"  # bump if index structure changes

# Course/year nesting:
# domains/Courses/<Course>/YEAR_X/SEM_Y/<Lesson>
COURSES_ROOT = os.path.join(DATA_ROOT, "Courses")

if SentenceTransformer is None:
    raise RuntimeError("sentence-transformers is required")
if faiss is None:
    raise RuntimeError("faiss is required")

embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2", device=EMBEDDING_DEVICE)

PROMPT_LANG = (
    "Απάντησε πάντα στα Ελληνικά, με φυσική και κατανοητή γλώσσα. "
    "Αν η ερώτηση είναι σε άλλη γλώσσα, μετέφρασέ την συνοπτικά και απάντησε στα Ελληνικά."
)

# =========================
# In-memory indexes
# =========================
# domain_data:
# { domainId: { "chunks": [...], "index": faiss.IndexFlatIP, "meta": [...],
#              "embeddings": np.ndarray, "folder": "...", "kind": "lesson|year",
#              "children": [...]} }
domain_data = {}

# domainId -> folder path
domain_folders = {}

# leaf lesson domain -> year aggregate domain (YEAR__Course__YEAR_X)
lesson_to_year = {}

# ✅ NEW: year alias -> year aggregate domain
# e.g. "Courses__Psychology__YEAR_1" -> "YEAR__Psychology__YEAR_1"
year_alias_to_year = {}


# =========================
# Helpers
# =========================
def require_api_key():
    key = request.headers.get("x-api-key")
    if not API_KEY or API_KEY == "changeme":
        print("⚠️ WARNING: API_KEY is default. Set API_KEY env var before public exposure!")
    elif key != API_KEY:
        abort(401, description="Unauthorized")


def sliding_word_chunks(words, max_tokens, overlap_tokens):
    chunks = []
    cur = []
    cur_len = 0
    for w in words:
        cur.append(w)
        cur_len += 1
        if cur_len >= max_tokens:
            chunks.append(" ".join(cur))
            if overlap_tokens > 0:
                cur = cur[-overlap_tokens:]
                cur_len = len(cur)
            else:
                cur = []
                cur_len = 0
    if cur:
        chunks.append(" ".join(cur))
    return chunks


def normalize_text(s):
    if not s:
        return ""
    s = unicodedata.normalize("NFKC", s)
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def safe_read_pdf_text(path):
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            texts = []
            for page in reader.pages:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t:
                    texts.append(t)
            return "\n".join(texts)
    except Exception as e:
        print(f"Error reading PDF {path}: {e}")
        return ""


def safe_read_docx_text(path):
    if docx is None:
        return ""
    try:
        d = docx.Document(path)
        return "\n".join([p.text for p in d.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX {path}: {e}")
        return ""


def safe_read_pptx_text(path):
    if pptx is None:
        return ""
    try:
        pres = pptx.Presentation(path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"Error reading PPTX {path}: {e}")
        return ""


def safe_read_xls_text(path):
    if xlrd is None and openpyxl is None:
        return ""
    try:
        text = []
        if path.lower().endswith(".xls") and xlrd is not None:
            wb = xlrd.open_workbook(path)
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    row_vals = [str(sheet.cell(row_idx, col_idx).value) for col_idx in range(sheet.ncols)]
                    text.append(" ".join(row_vals))
        elif path.lower().endswith((".xlsx", ".xlsm")) and openpyxl is not None:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(c) for c in row if c is not None]
                    if row_vals:
                        text.append(" ".join(row_vals))
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading XLS/XLSX {path}: {e}")
        return ""


def read_file_text(path):
    lower = path.lower()
    if lower.endswith(".pdf"):
        return safe_read_pdf_text(path)
    elif lower.endswith(".docx"):
        return safe_read_docx_text(path)
    elif lower.endswith(".pptx"):
        return safe_read_pptx_text(path)
    elif lower.endswith((".xls", ".xlsx", ".xlsm")):
        return safe_read_xls_text(path)
    else:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file {path}: {e}")
            return ""


def tokenize_rough(text):
    # very rough: count words as tokens
    return len(text.split())


def embed_texts(texts):
    if not texts:
        return np.zeros((0, 384), dtype=np.float32)
    all_vecs = []
    for i in range(0, len(texts), EMBED_BATCH_SIZE):
        batch = texts[i:i + EMBED_BATCH_SIZE]
        emb = embedding_model.encode(batch, convert_to_numpy=True, normalize_embeddings=True)
        all_vecs.append(emb.astype(np.float32))
    return np.vstack(all_vecs)


def build_index_for_domain(domain_id, folder, kind="lesson", children=None):
    """
    Build a FAISS index for all files in `folder`.
    kind: "lesson" or "year"
    children: list of child lesson domain ids (for year aggregates)
    """
    print(f"Building index for domain={domain_id}, folder={folder}, kind={kind}")
    chunks = []
    meta = []

    # Collect chunks
    for root, dirs, files in os.walk(folder):
        for fname in files:
            fpath = os.path.join(root, fname)
            rel = os.path.relpath(fpath, folder)
            lower = fname.lower()

            if lower.endswith((".pdf", ".txt", ".docx", ".pptx", ".xls", ".xlsx", ".xlsm")):
                text = read_file_text(fpath)
                text = normalize_text(text)
                if not text:
                    continue
                words = text.split()
                if len(words) < MIN_WORDS_TO_INDEX:
                    continue

                # chunk
                w_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
                page_hint = None
                # quick heuristic for PDFs pages
                if lower.endswith(".pdf"):
                    try:
                        reader = PyPDF2.PdfReader(fpath)
                        # naive per-page text: we re-use safe_read_pdf_text above, but we only have full text here;
                        # so we just put page_hint=None or 1
                        page_hint = 1
                    except Exception:
                        page_hint = None

                for ch in w_chunks:
                    chunks.append(ch)
                    meta.append(
                        {
                            "file": fname,
                            "rel": rel,
                            "path": fpath,
                            "page": page_hint,
                        }
                    )

    if not chunks:
        print(f"No chunks for domain {domain_id}")
        embeddings = np.zeros((0, 384), dtype=np.float32)
        index = faiss.IndexFlatIP(384)
    else:
        embeddings = embed_texts(chunks)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)

    domain_data[domain_id] = {
        "chunks": chunks,
        "embeddings": embeddings,
        "index": index,
        "meta": meta,
        "folder": folder,
        "kind": kind,
        "children": children or [],
    }

    domain_folders[domain_id] = folder

    print(
        f"Built index for domain={domain_id}, kind={kind}. "
        f"Chunks={len(chunks)}, folder={folder}, children={children}"
    )
    return domain_data[domain_id]


def _cache_key_for_domain(domain_id, folder, kind, children):
    h = hashlib.sha256()
    h.update(domain_id.encode("utf-8", errors="ignore"))
    h.update(folder.encode("utf-8", errors="ignore"))
    h.update(kind.encode("utf-8", errors="ignore"))
    if children:
        for c in sorted(children):
            h.update(c.encode("utf-8", errors="ignore"))
    h.update(INDEX_CACHE_VERSION.encode("utf-8"))
    return h.hexdigest()


def save_index_to_disk(domain_id):
    if not USE_INDEX_CACHE:
        return
    info = domain_data.get(domain_id)
    if not info:
        return
    folder = info["folder"]
    kind = info["kind"]
    children = info.get("children") or []

    cache_key = _cache_key_for_domain(domain_id, folder, kind, children)
    base = os.path.join(INDEX_CACHE_DIR, cache_key)
    os.makedirs(base, exist_ok=True)

    # Save meta
    meta_path = os.path.join(base, "meta.json")
    data = {
        "domain_id": domain_id,
        "folder": folder,
        "kind": kind,
        "children": children,
        "chunks": info["chunks"],
        "meta": info["meta"],
    }
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False)

    # Save embeddings
    emb_path = os.path.join(base, "embeddings.npy")
    np.save(emb_path, info["embeddings"])

    # Save FAISS index
    index_path = os.path.join(base, "index.faiss")
    faiss.write_index(info["index"], index_path)

    print(f"Saved index cache for domain={domain_id} -> {base}")


def load_index_from_disk(domain_id, folder, kind, children):
    if not USE_INDEX_CACHE:
        return None
    cache_key = _cache_key_for_domain(domain_id, folder, kind, children)
    base = os.path.join(INDEX_CACHE_DIR, cache_key)
    meta_path = os.path.join(base, "meta.json")
    emb_path = os.path.join(base, "embeddings.npy")
    index_path = os.path.join(base, "index.faiss")
    if not (os.path.isfile(meta_path) and os.path.isfile(emb_path) and os.path.isfile(index_path)):
        return None
    try:
        with open(meta_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        embeddings = np.load(emb_path)
        index = faiss.read_index(index_path)
    except Exception as e:
        print(f"Failed to load index cache for {domain_id}: {e}")
        return None

    domain_data[domain_id] = {
        "chunks": data["chunks"],
        "embeddings": embeddings,
        "index": index,
        "meta": data["meta"],
        "folder": data["folder"],
        "kind": data["kind"],
        "children": data.get("children") or [],
    }
    domain_folders[domain_id] = data["folder"]
    print(f"Loaded index cache for domain={domain_id} from {base}")
    return domain_data[domain_id]


def ensure_domain_index(domain_id, folder, kind="lesson", children=None):
    if domain_id in domain_data:
        return domain_data[domain_id]
    # try cache
    if load_index_from_disk(domain_id, folder, kind, children):
        return domain_data[domain_id]
    # build fresh
    info = build_index_for_domain(domain_id, folder, kind, children)
    save_index_to_disk(domain_id)
    return info


def discover_flat_domains():
    """
    Legacy flat domains: DATA_ROOT/<domainId> with PDFs, etc.
    """
    if not os.path.isdir(DATA_ROOT):
        return
    for name in os.listdir(DATA_ROOT):
        if name == "Courses":
            continue
        path = os.path.join(DATA_ROOT, name)
        if os.path.isdir(path):
            domain_folders[name] = path
            # index lazily on demand
            print(f"Discovered flat domain: {name} -> {path}")


def slugify_course_name(name):
    # e.g. "Psychology" or "Ψυχολογία" -> used only in YEAR_ aggregate id
    # keep original for now; domain ids will be constructed as:
    # YEAR__<CourseName>__YEAR_X
    name = name.replace(" ", "_")
    return name


def discover_courses_hierarchy():
    """
    Discover nested Courses/<Course>/YEAR_X/SEM_Y/<Lesson>.
    For each leaf lesson, register a domain id:
      Courses__<Course>__YEAR_X__SEM_Y__<Lesson>
    Also build a year aggregate domain:
      YEAR__<Course>__YEAR_X
    """
    if not os.path.isdir(COURSES_ROOT):
        print(f"No Courses root at {COURSES_ROOT}")
        return
    for course_name in os.listdir(COURSES_ROOT):
        course_path = os.path.join(COURSES_ROOT, course_name)
        if not os.path.isdir(course_path):
            continue

        for year_name in os.listdir(course_path):
            if not re.match(r"YEAR_\d+", year_name, re.IGNORECASE):
                continue
            year_path = os.path.join(course_path, year_name)
            if not os.path.isdir(year_path):
                continue

            # Aggregate year domain id
            course_slug = slugify_course_name(course_name)
            year_dom = f"YEAR__{course_slug}__{year_name}"

            lesson_domains = []

            for sem_name in os.listdir(year_path):
                sem_path = os.path.join(year_path, sem_name)
                if not os.path.isdir(sem_path):
                    continue

                for lesson_name in os.listdir(sem_path):
                    lesson_path = os.path.join(sem_path, lesson_name)
                    if not os.path.isdir(lesson_path):
                        continue

                    lesson_dom = f"Courses__{course_name}__{year_name}__{sem_name}__{lesson_name}"
                    domain_folders[lesson_dom] = lesson_path
                    lesson_domains.append(lesson_dom)
                    # mapping lesson -> year aggregate
                    lesson_to_year[lesson_dom] = year_dom

                    print(f"Discovered lesson domain: {lesson_dom} -> {lesson_path}")

            if lesson_domains:
                # The folder for the year aggregate domain is the YEAR_X folder
                domain_folders[year_dom] = year_path
                domain_data[year_dom] = {
                    "folder": year_path,
                    "kind": "year",
                    "children": lesson_domains,
                    "chunks": [],
                    "embeddings": np.zeros((0, 384), dtype=np.float32),
                    "index": faiss.IndexFlatIP(384),
                    "meta": [],
                }
                print(
                    f"Registered year aggregate domain: {year_dom} -> {year_path}, "
                    f"children={len(lesson_domains)}"
                )

                # ✅ also register a "flat-style" alias so frontend can just send:
                #    Courses__<Course>__YEAR_X
                alias_key = f"Courses__{course_name}__{year_name}"
                alias_key = alias_key.replace(" ", "_")
                year_alias_to_year[alias_key] = year_dom
                print(f"Alias registered: {alias_key} -> {year_dom}")


def init_domains():
    discover_flat_domains()
    discover_courses_hierarchy()


init_domains()


def _require_cache_admin():
    key = request.headers.get("x-cache-admin")
    if CACHE_ADMIN_KEY and key == CACHE_ADMIN_KEY:
        return
    abort(403, description="Cache admin key required")


# =========================
# Retrieval
# =========================
def retrieve_chunks(domain_id, query, top_k=R_TOP_K, alt_queries=None):
    """
    Retrieve top_k chunks for a query (and alternate queries) from domain.
    Returns: list of (chunk_text, score, meta)
    """
    if domain_id not in domain_folders:
        raise ValueError(f"Unknown domain: {domain_id}")
    info = domain_data.get(domain_id)
    if not info or "index" not in info:
        folder = domain_folders[domain_id]
        kind = info.get("kind", "lesson") if info else "lesson"
        children = info.get("children") if info else None
        info = ensure_domain_index(domain_id, folder, kind, children)

    index = info["index"]
    embeddings = info["embeddings"]
    chunks = info["chunks"]
    meta = info["meta"]

    if embeddings.shape[0] == 0:
        return []

    queries = [query] + (alt_queries or [])
    q_emb = embed_texts(queries)
    D, I = index.search(q_emb, top_k)

    results = []
    used = set()
    for row in range(D.shape[0]):
        for col in range(D.shape[1]):
            idx = int(I[row, col])
            if idx < 0 or idx >= len(chunks):
                continue
            if idx in used:
                continue
            score = float(D[row, col])
            if score < MIN_SIM_THRESHOLD:
                continue
            used.add(idx)
            results.append((chunks[idx], score, meta[idx]))

    results.sort(key=lambda x: x[1], reverse=True)
    return results[:top_k]


def build_prompt_from_chunks(question, retrieved, chat_messages=None):
    """
    Build the final prompt given retrieved chunks and a question.
    """
    context_parts = []
    sources = []

    for i, (chunk, score, m) in enumerate(retrieved, start=1):
        file = m.get("file")
        page = m.get("page")
        snippet = chunk
        if len(snippet) > 240:
            snippet = snippet[:240] + "…"

        sources.append(
            {
                "label": i,
                "file": file,
                "page_start": page,
                "page_end": page,
                "snippet": snippet,
            }
        )

        context_parts.append(f"[{i}] ({file}, σελίδα {page}) {chunk}")

    history_str = ""
    if chat_messages:
        # we only put a short window of previous user/assistant messages
        history_items = []
        for msg in chat_messages[-4:]:
            role = msg.get("role")
            content = msg.get("content") or ""
            if role == "user":
                history_items.append(f"Χρήστης: {content}")
            elif role == "assistant":
                history_items.append(f"Βοηθός: {content}")
        if history_items:
            history_str = "\n\nΙστορικό συνομιλίας:\n" + "\n".join(history_items)

    context_text = "\n\n".join(context_parts) if context_parts else ""
    instructions = (
        f"{PROMPT_LANG}\n"
        "Χρησιμοποίησε ΜΟΝΟ τις παρακάτω πληροφορίες από τα έγγραφα για να απαντήσεις. "
        "Αν οι πληροφορίες δεν επαρκούν, πες ξεκάθαρα ότι δεν είσαι βέβαιος.\n"
    )

    full_prompt = (
        instructions
        + history_str
        + ("\n\nΠληροφορίες από έγγραφα:\n" + context_text if context_text else "")
        + f"\n\nΕρώτηση: {question}\n\nΑπάντηση:"
    )

    return full_prompt, sources


def call_vllm(prompt, max_tokens=MAX_TOKENS):
    """
    Call vLLM chat endpoint with token budgeting.
    """
    # crude estimate: 1 token ~ 4 chars
    est_prompt_tokens = max(1, int(len(prompt) / 4))
    max_context = VLLM_MAX_CONTEXT
    safety = VLLM_SAFETY_MARGIN

    avail_for_completion = max_context - est_prompt_tokens - safety
    avail_for_completion = max(avail_for_completion, VLLM_MIN_COMPLETION)
    avail_for_completion = min(avail_for_completion, VLLM_MAX_COMPLETION_CAP)

    payload = {
        "model": VLLM_MODEL,
        "messages": [
            {"role": "system", "content": "You are a helpful AI tutor that answers in Greek."},
            {"role": "user", "content": prompt},
        ],
        "temperature": TEMPERATURE,
        "top_p": TOP_P,
        "max_tokens": avail_for_completion,
    }

    try:
        resp = requests.post(VLLM_URL, json=payload, timeout=VLLM_TIMEOUT)
    except Exception as e:
        print(f"Error calling vLLM: {e}")
        raise RuntimeError("vLLM call failed")

    if resp.status_code != 200:
        print(f"vLLM error {resp.status_code}: {resp.text[:500]}")
        raise RuntimeError(f"vLLM HTTP {resp.status_code}")

    data = resp.json()
    try:
        content = data["choices"][0]["message"]["content"]
    except Exception:
        content = ""
    return content.strip() or "Δεν μπόρεσα να παραγάγω απάντηση."


def answer_question(domain_id, question, chat_messages=None):
    """
    Main QA pipeline for a given domain.
    """
    alt_queries = []
    if NUM_ALTERNATE_QUERIES > 0:
        alt_queries = [question.upper()] if NUM_ALTERNATE_QUERIES >= 1 else []

    retrieved = retrieve_chunks(domain_id, question, top_k=R_TOP_K, alt_queries=alt_queries)
    if not retrieved:
        if STRICT_MODE:
            return (
                "Δεν βρήκα σχετικές πληροφορίες στα έγγραφα αυτού του μαθήματος / έτους για να απαντήσω με σιγουριά.",
                [],
            )
        else:
            # no context, still call model
            prompt = (
                f"{PROMPT_LANG}\n"
                "Δυστυχώς δεν βρήκα σχετικές πληροφορίες στα έγγραφα, αλλά προσπάθησε να απαντήσεις "
                "με βάση τις γενικές σου γνώσεις.\n\n"
                f"Ερώτηση: {question}\n\nΑπάντηση:"
            )
            ans = call_vllm(prompt)
            return ans, []

    prompt, sources = build_prompt_from_chunks(question, retrieved, chat_messages=chat_messages)

    if EXTRACTIVE_ONLY:
        # in extractive mode, don't call model—just return a concatenation of top chunks
        joined = "\n\n".join([c for (c, _, _) in retrieved])
        return joined, sources

    ans = call_vllm(prompt)
    return ans, sources


# =========================
# Teacher answers (per-domain PDF)
# =========================
def load_teacher_answers_for_domain(domain: str):
    """
    Load and cache teacher answers for a given domain.

    Expected file name inside the domain folder:
      {domain}_TEACHER_ANSWERS.pdf

    Example:
      domain = "Courses__Psychology__YEAR_2__SEM_3__ΨΥΧΓΕ3"
      file   = "Courses__Psychology__YEAR_2__SEM_3__ΨΥΧΓΕ3_TEACHER_ANSWERS.pdf"
    """
    if not domain:
        return None

    # Cache hit
    if domain in teacher_answers_cache:
        return teacher_answers_cache[domain]

    folder = _resolve_domain_folder(domain)
    if not folder:
        teacher_answers_cache[domain] = None
        return None

    teacher_fname = f"{domain}_TEACHER_ANSWERS.pdf"
    teacher_path = os.path.join(folder, teacher_fname)
    if not os.path.isfile(teacher_path):
        teacher_answers_cache[domain] = None
        return None

    questions = []
    answers = []
    pages = []

    try:
        reader = PyPDF2.PdfReader(teacher_path)
    except Exception as e:
        print(f"⚠️ Could not read teacher PDF for domain {domain}: {e}")
        teacher_answers_cache[domain] = None
        return None

    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""

        if not text.strip():
            continue

        # Simple parsing based on the expected pattern
        q_key = "Ερώτηση:"
        a_key = "Απάντηση καθηγητή:"
        d_key = "DocId:"

        q_idx = text.find(q_key)
        a_idx = text.find(a_key)

        if q_idx == -1 or a_idx == -1 or a_idx <= q_idx:
            continue

        q_text = text[q_idx + len(q_key) : a_idx].strip()

        d_idx = text.find(d_key, a_idx)
        if d_idx != -1:
            a_text = text[a_idx + len(a_key) : d_idx].strip()
        else:
            a_text = text[a_idx + len(a_key) :].strip()

        if q_text and a_text:
            questions.append(q_text)
            answers.append(a_text)
            pages.append(i + 1)  # 1-based page index

    if not questions:
        teacher_answers_cache[domain] = None
        return None

    try:
        emb = embedding_model.encode(
            questions, convert_to_numpy=True, normalize_embeddings=True
        ).astype(np.float32)
    except Exception as e:
        print(f"⚠️ Failed to embed teacher questions for domain {domain}: {e}")
        teacher_answers_cache[domain] = None
        return None

    teacher_answers_cache[domain] = {
        "questions": questions,
        "answers": answers,
        "embeddings": emb,
        "pages": pages,
        "file": teacher_fname,
    }
    print(f"📘 Loaded {len(questions)} teacher Q&A pairs for domain '{domain}'.")
    return teacher_answers_cache[domain]


def maybe_answer_from_teacher(domain: str, question: str):
    """
    Try to answer a question directly from the teacher answers PDF for this domain.
    Returns (answer_text, sources) or (None, None) if no good match.
    """
    if not question or not question.strip():
        return None, None

    info = load_teacher_answers_for_domain(domain)
    if not info:
        return None, None

    q_emb = embedding_model.encode(
        [question], convert_to_numpy=True, normalize_embeddings=True
    ).astype(np.float32)[0]

    emb_mat = info.get("embeddings")
    if emb_mat is None or not len(emb_mat):
        return None, None

    # cosine similarity (embeddings are normalized)
    sims = emb_mat @ q_emb
    best_idx = int(np.argmax(sims))
    best_sim = float(sims[best_idx])

    if best_sim < TEACHER_MIN_SIM:
        return None, None

    answer_text = info["answers"][best_idx]
    page = info["pages"][best_idx] if info.get("pages") else None
    teacher_file = info.get("file")

    snippet = answer_text
    if len(snippet) > 240:
        snippet = snippet[:240] + "…"

    sources = [
        {
            "label": 1,
            "file": teacher_file,
            "page_start": page,
            "page_end": page,
            "snippet": snippet,
        }
    ]

    return answer_text, sources


# =========================
# Domain folder resolving
# =========================
def _resolve_domain_folder(domain: str):
    if domain in domain_folders:
        return domain_folders[domain]
    d2 = next((k for k in domain_folders.keys() if k.lower() == (domain or "").lower()), None)
    if d2:
        return domain_folders[d2]
    return None


@app.route("/files/<domain>/<path:filename>", methods=["GET"])
def serve_file(domain, filename):
    if not filename.lower().endswith(".pdf"):
        abort(403)
    folder = _resolve_domain_folder(domain)
    if not folder or not os.path.isdir(folder):
        abort(404)
    safe_path = safe_join(folder, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)
    return send_file(safe_path, as_attachment=False)


@app.route("/domains", methods=["GET"])
def list_domains():
    # Return a list of known domain IDs and some meta
    items = []
    for dom, folder in domain_folders.items():
        info = domain_data.get(dom) or {}
        items.append(
            {
                "domain": dom,
                "folder": folder,
                "kind": info.get("kind", "lesson"),
                "children": info.get("children", []),
            }
        )
    return jsonify(items)


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "model": MODEL_NAME})


@app.route("/ask", methods=["POST"])
def ask():
    require_api_key()
    data = request.get_json(force=True) or {}

    question = (data.get("question") or "").strip()
    domain_in = (data.get("domain") or "").strip()
    scope = (data.get("scope") or "lesson").strip().lower()  # lesson | year
    chat_messages = data.get("messages") or []

    if not question:
        return jsonify({"answer": "Δεν δόθηκε ερώτηση.", "sources": []})

    # match domain case-insensitively
    matched_domain = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), domain_in)

    # ✅ scope escalation (supports BOTH lesson id and year-alias id)
    if scope == "year":
        # 1) if frontend sends "Courses__Course__YEAR_X", map to year aggregate
        alias_key = domain_in.replace(" ", "_")
        if alias_key in year_alias_to_year and year_alias_to_year[alias_key] in domain_data:
            matched_domain = year_alias_to_year[alias_key]
        else:
            # 2) if it's a lesson id, map to its year aggregate
            yd = lesson_to_year.get(matched_domain)
            if yd and yd in domain_data:
                matched_domain = yd
            else:
                # 3) if frontend already sent YEAR__..., keep it
                yd2 = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), None)
                if yd2 and domain_data.get(yd2, {}).get("kind") == "year":
                    matched_domain = yd2

    print(f"Question for domain [{matched_domain}] (scope={scope}): {question}")

    # 1) Try teacher answers first (only if there is a teacher PDF for this domain)
    teacher_answer, teacher_sources = maybe_answer_from_teacher(matched_domain, question)
    if teacher_answer is not None:
        # Identifier for frontend: teacherAnswer = true
        return jsonify(
            {
                "answer": teacher_answer,
                "sources": teacher_sources or [],
                "resolvedDomain": matched_domain,
                "teacherAnswer": True,
            }
        )

    # 2) Fallback to existing retrieval + vLLM pipeline (unchanged)
    answer, sources = answer_question(matched_domain, question, chat_messages=chat_messages)

    return jsonify({"answer": answer, "sources": sources, "resolvedDomain": matched_domain})


# =========================
# Cache admin (optional)
# =========================
@app.route("/cache/clear", methods=["POST"])
def cache_clear():
    _require_cache_admin()
    data = request.get_json(force=True, silent=True) or {}
    domain = data.get("domain")
    reload_after = bool(data.get("reload", False))

    summary = clear_index_cache(domain=domain, reload_after=reload_after)
    return jsonify(summary)


def clear_index_cache(domain=None, reload_after=False):
    """
    Clear index caches from disk and/or memory.
    If domain is None, clear all.
    """
    removed_domains = 0
    removed_files = 0

    if domain:
        # specific domain
        info = domain_data.pop(domain, None)
        domain_folders.pop(domain, None)
        cache_key = None
        if info:
            cache_key = _cache_key_for_domain(domain, info["folder"], info.get("kind", "lesson"), info.get("children"))
        else:
            # try to guess folder if missing
            folder = domain_folders.get(domain, "")
            cache_key = _cache_key_for_domain(domain, folder, "lesson", None)
        if cache_key:
            base = os.path.join(INDEX_CACHE_DIR, cache_key)
            if os.path.isdir(base):
                for fname in os.listdir(base):
                    fp = os.path.join(base, fname)
                    try:
                        os.remove(fp)
                        removed_files += 1
                    except Exception:
                        pass
                try:
                    shutil.rmtree(base, ignore_errors=True)
                except Exception:
                    pass
                removed_domains += 1
    else:
        # all domains
        domain_data.clear()
        domain_folders.clear()

        for t in os.listdir(INDEX_CACHE_DIR):
            base = os.path.join(INDEX_CACHE_DIR, t)
            if os.path.isdir(base):
                for fname in os.listdir(base):
                    fp = os.path.join(base, fname)
                    try:
                        os.remove(fp)
                        removed_files += 1
                    except Exception:
                        pass
            try:
                shutil.rmtree(base, ignore_errors=True)
            except Exception:
                pass

    os.makedirs(INDEX_CACHE_DIR, exist_ok=True)
    return {
        "domain": domain or "*",
        "removed_domains": removed_domains,
        "removed_files": removed_files,
        "reloaded": bool(reload_after),
    }


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.getenv("PORT", "5000")), debug=False)
