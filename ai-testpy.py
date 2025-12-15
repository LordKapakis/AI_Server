# server.py
from flask import Flask, request, jsonify, abort, send_from_directory, send_file
import os
import io
import re
import PyPDF2
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import requests
from flask_cors import CORS
from werkzeug.utils import safe_join  # correct import

# NEW IMPORTS FOR DOCX & PPTX
from docx import Document
from pptx import Presentation

app = Flask(__name__)
CORS(app)

# === Core Config ===
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DEFAULT_BASE = os.path.join(SCRIPT_DIR, "domains")
BASE_FOLDER = os.getenv("BASE_FOLDER", DEFAULT_BASE)
os.makedirs(BASE_FOLDER, exist_ok=True)

OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
MODEL_NAME = os.getenv("MODEL_NAME", "llama3.1")
API_KEY = os.getenv("API_KEY", "changeme")  # set a strong secret before exposing

NUM_CTX = int(os.getenv("NUM_CTX", "4096"))
NUM_PREDICT = int(os.getenv("NUM_PREDICT", "400"))
TEMPERATURE = float(os.getenv("TEMPERATURE", "0.2"))
TOP_P = float(os.getenv("TOP_P", "0.9"))

# === Retrieval / Index Tuning ===
CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "180"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "30"))
R_TOP_K = int(os.getenv("R_TOP_K", "4"))
NUM_ALTERNATE_QUERIES = int(os.getenv("NUM_ALTERNATE_QUERIES", "1"))
MIN_SIM_THRESHOLD = float(os.getenv("MIN_SIM_THRESHOLD", "0.30"))  # used in post-filter

# Be strict about missing citations
STRICT_MODE = os.getenv("STRICT_MODE", "true").lower() in ("1", "true", "yes")
EXTRACTIVE_ONLY = os.getenv("EXTRACTIVE_ONLY", "false").lower() in ("1", "true", "yes")

# Multilingual embedding model (handles Greek well)
embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")

# Greek output by default
PROMPT_LANG = (
    "Απάντησε πάντα στα Ελληνικά, με φυσική και κατανοητή γλώσσα. "
    "Αν η ερώτηση είναι σε άλλη γλώσσα, μετέφρασέ την συνοπτικά και απάντησε στα Ελληνικά."
)

# domain_data = { domain: { "chunks": [...], "index": faiss.IndexFlatIP, "meta": [...] } }
domain_data = {}


# -----------------------------
# Helpers: chunking and loading
# -----------------------------
def sliding_word_chunks(words, size=CHUNK_TOKENS, overlap=CHUNK_OVERLAP):
    """Return overlapping word chunks to improve recall."""
    if size <= 0:
        return []
    step = max(1, size - overlap)
    out = []
    for i in range(0, len(words), step):
        piece = words[i: i + size]
        if not piece:
            break
        out.append(" ".join(piece))
        if i + size >= len(words):
            break
    return out


def pdf_to_chunks_with_meta(pdf_path, chunk_size=CHUNK_TOKENS, overlap=CHUNK_OVERLAP):
    """
    Returns: (chunks, metas)
    metas: list of dicts: {"file": <name>, "page_start": int, "page_end": int, "snippet": str}
    Robust mapping: builds a word_to_page list so chunks always map to correct pages,
    even when some PDF pages extract as empty.
    """
    chunks, metas = [], []
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        per_page_text = []
        for pageno in range(len(reader.pages)):
            t = reader.pages[pageno].extract_text() or ""
            per_page_text.append(" ".join(t.split()))  # light normalization

        # Flat words array + word->page lookup
        words = []
        word_to_page = []  # same length as words; stores 0-based page number
        for pageno, text in enumerate(per_page_text):
            if text:
                w = text.split()
                words.extend(w)
                word_to_page.extend([pageno] * len(w))
            else:
                # empty page: contributes no words; mapping stays consistent
                pass

        if not words:
            # If the PDF has no extractable text at all, skip it
            return [], []

        step = max(1, chunk_size - overlap)
        i = 0
        basename = os.path.basename(pdf_path)

        while i < len(words):
            j = min(i + chunk_size, len(words))
            piece_words = words[i:j]
            if not piece_words:
                break

            pages_in_slice = word_to_page[i:j]
            if pages_in_slice:
                page_start0 = min(pages_in_slice)
                page_end0 = max(pages_in_slice)
            else:
                page_start0 = page_end0 = 0  # defensive

            text_chunk = " ".join(piece_words)

            chunks.append(text_chunk)
            metas.append({
                "file": basename,
                "page_start": int(page_start0 + 1),  # 1-based for UI
                "page_end": int(page_end0 + 1),
                "snippet": (text_chunk[:240] + "…") if len(text_chunk) > 240 else text_chunk
            })

            if j >= len(words):
                break
            i += step

    except Exception as e:
        print(f"❌ Could not read {pdf_path}: {e}")
    return chunks, metas


# NEW: DOCX helper
def docx_to_text(docx_path):
    """Extract plain text from a .docx file."""
    doc = Document(docx_path)
    parts = []

    # Paragraph text
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)

    # Table text (optional but useful)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)

    return "\n".join(parts)


# NEW: PPTX helper
def pptx_to_text(pptx_path):
    """Extract plain text from a .pptx file."""
    pres = Presentation(pptx_path)
    parts = []

    for slide in pres.slides:
        for shape in slide.shapes:
            # Regular text frames
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text
                if text:
                    parts.append(text)

            # Tables
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            parts.append(text)

    return "\n".join(parts)


def build_domain_index(domain, folder):
    chunks, sources_meta = [], []
    files_found = 0

    for fname in os.listdir(folder):
        # EXTENDED: now supports .pdf, .txt, .docx, .pptx
        if fname.lower().endswith((".pdf", ".txt", ".docx", ".pptx")):
            files_found += 1
            path = os.path.join(folder, fname)
            print(f"📖 Loading {domain}: {path}")

            # --- PDF ---
            if fname.lower().endswith(".pdf"):
                new_chunks, metas = pdf_to_chunks_with_meta(path, CHUNK_TOKENS, CHUNK_OVERLAP)
                chunks.extend(new_chunks)
                sources_meta.extend(metas)

            # --- TXT ---
            elif fname.lower().endswith(".txt"):
                try:
                    with open(path, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()
                except Exception as e:
                    print(f"❌ Could not read TXT {path}: {e}")
                    continue

                words = (text or "").split()
                new_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
                chunks.extend(new_chunks)
                for ch in new_chunks:
                    sources_meta.append({
                        "file": os.path.basename(path),
                        "page_start": None,
                        "page_end": None,
                        "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
                    })

            # --- DOCX ---
            elif fname.lower().endswith(".docx"):
                try:
                    text = docx_to_text(path)
                except Exception as e:
                    print(f"❌ Could not read DOCX {path}: {e}")
                    continue

                words = (text or "").split()
                new_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
                chunks.extend(new_chunks)
                for ch in new_chunks:
                    sources_meta.append({
                        "file": os.path.basename(path),
                        "page_start": None,
                        "page_end": None,
                        "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
                    })

            # --- PPTX ---
            elif fname.lower().endswith(".pptx"):
                try:
                    text = pptx_to_text(path)
                except Exception as e:
                    print(f"❌ Could not read PPTX {path}: {e}")
                    continue

                words = (text or "").split()
                new_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
                chunks.extend(new_chunks)
                for ch in new_chunks:
                    sources_meta.append({
                        "file": os.path.basename(path),
                        "page_start": None,
                        "page_end": None,
                        "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
                    })

    if files_found == 0 or not chunks:
        print(f"⚠️ No usable data for '{domain}' in {folder}")
        return

    embeddings = embedding_model.encode(chunks, convert_to_numpy=True, normalize_embeddings=True)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)  # cosine on normalized vectors via dot on normalized
    index.add(embeddings)

    domain_data[domain] = {"chunks": chunks, "index": index, "meta": sources_meta}
    print(f"✅ Indexed {len(chunks)} chunks for '{domain}'")


# -----------------------------
# LLM call
# -----------------------------
def query_model(prompt: str, model: str = MODEL_NAME) -> str:
    # prepend the language rule
    full_prompt = f"{PROMPT_LANG}\n\n{prompt}"
    payload = {
        "model": model,
        "prompt": full_prompt,
        "stream": False,
        "options": {
            "num_ctx": NUM_CTX,
            "num_predict": NUM_PREDICT,
            "temperature": TEMPERATURE,
            "top_p": TOP_P,
        },
    }
    resp = requests.post(f"{OLLAMA_URL}/api/generate", json=payload, timeout=300)
    resp.raise_for_status()
    data = resp.json()
    return data.get("response", "").strip()


def generate_alternate_queries(user_q: str, n: int = NUM_ALTERNATE_QUERIES):
    """Generate a few paraphrases to improve recall without broadening the topic or changing language."""
    if n <= 0:
        return [user_q]

    prompt = f"""
Παράφρασε την ερώτηση του χρήστη σε έως {n} εναλλακτικές διατυπώσεις για ΑΝΑΖΗΤΗΣΗ.
Απαιτήσεις:
- Διατήρησε το ίδιο νόημα, μην προσθέτεις/αφαιρείς υποθέματα.
- Χρησιμοποίησε την ΙΔΙΑ γλώσσα με την αρχική ερώτηση.
- Απόφυγε υπερ-γενικεύσεις και αλλαγές πεδίου.
- Κάθε παραλλαγή σε ξεχωριστή γραμμή, χωρίς αρίθμηση/περιττό κείμενο.

Ερώτηση: "{user_q}"
"""
    try:
        text = query_model(prompt)
    except Exception as e:
        print("⚠️ Alternate query generation failed:", e)
        return [user_q]

    alts = [line.strip("- ").strip() for line in text.splitlines() if line.strip()]
    all_qs = [user_q] + alts[:n]

    seen = set()
    uniq = []
    for q in all_qs:
        key = q.lower()
        if key not in seen:
            seen.add(key)
            uniq.append(q)
    return uniq


def build_labeled_context(chunks):
    labeled = []
    for i, ch in enumerate(chunks, start=1):
        labeled.append(f"[{i}] {ch}")
    return "\n\n".join(labeled), list(range(1, len(chunks) + 1))


# --- citations safety net ---
def ensure_citations_from_quotes(final_answer: str, quotes: str, allowed_labels):
    """If final_answer has no [n], append the unique labels seen in QUOTES, filtered to allowed."""
    if re.search(r"\[\d+\]", final_answer):
        return final_answer  # already cited

    quoted_labels = [int(n) for n in re.findall(r"\[(\d+)\]", quotes)]
    if not quoted_labels:
        return final_answer

    allowed_set = set(int(x) for x in allowed_labels)
    seen = []
    for n in quoted_labels:
        if n in allowed_set and n not in seen:
            seen.append(n)

    if not seen:
        return final_answer

    tail = "".join(f"[{n}]" for n in seen)
    return (final_answer.rstrip() + " " + tail).strip()


# -----------------------------
# Retrieval + Answer
# -----------------------------
def answer_question(domain, question, top_k=R_TOP_K):
    """
    Returns:
      final_answer: str
      sources_out: list[{"label": int, "file": str|None, "page_start": int|None,
                         "page_end": int|None, "snippet": str}]
    """
    if domain not in domain_data:
        return f"Domain '{domain}' not available.", []

    data = domain_data[domain]
    index = data["index"]
    chunks = data["chunks"]
    meta = data.get("meta")

    # Multi-query retrieval
    queries = generate_alternate_queries(question, NUM_ALTERNATE_QUERIES)

    scored = []
    for q in queries:
        q_emb = embedding_model.encode([q], convert_to_numpy=True, normalize_embeddings=True)
        S, I = index.search(q_emb, top_k)  # cosine similarity
        for sim, idx in zip(S[0].tolist(), I[0].tolist()):
            if 0 <= idx < len(chunks):
                scored.append((float(sim), int(idx)))

    if not scored:
        system_prompt = f"""
Είσαι βοηθός με ειδίκευση στο πεδίο «{domain}».
Δώσε μια σύντομη, υψηλού επιπέδου εξήγηση για το παρακάτω θέμα.
Αν δεν σχετίζεται με το πεδίο ή δεν υπάρχουν στοιχεία, απάντησε: "Δεν ξέρω."
"""
        prompt = f"{system_prompt}\n\nΕΡΩΤΗΣΗ: {question}\nΑΠΑΝΤΗΣΗ:"
        return query_model(prompt), []

    # Keep best similarity per chunk and apply optional threshold
    best_by_idx = {}
    for sim, idx in scored:
        best_by_idx[idx] = max(sim, best_by_idx.get(idx, -1.0))

    ranked = sorted(best_by_idx.items(), key=lambda x: x[1], reverse=True)

    # Apply similarity threshold (optional but useful)
    if MIN_SIM_THRESHOLD is not None:
        ranked = [(idx, sim) for idx, sim in ranked if sim >= MIN_SIM_THRESHOLD]

    ranked = ranked[:top_k]
    selected_idxs = [idx for idx, _ in ranked]

    # Retrieved texts and labeled context
    retrieved = [chunks[i] for i in selected_idxs]
    labeled_context = "\n\n".join(f"[{i+1}] {txt}" for i, txt in enumerate(retrieved))
    allowed_labels = [str(i + 1) for i in range(len(retrieved))]
    labels_hint = ", ".join(f"[{x}]" for x in allowed_labels)  # e.g. "[1], [2], [3]"

    # Build UI sources
    sources_out = []
    for n, idx in enumerate(selected_idxs, start=1):
        if meta and 0 <= idx < len(meta):
            m = meta[idx]
            file = m.get("file")
            p_start = m.get("page_start")
            p_end = m.get("page_end")
            snippet = m.get("snippet") or (retrieved[n - 1][:240] + "…"
                                           if len(retrieved[n - 1]) > 240 else retrieved[n - 1])
        else:
            file = None
            p_start = None
            p_end = None
            snippet = retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1]

        sources_out.append({
            "label": n,
            "file": file,
            "page_start": p_start,
            "page_end": p_end,
            "snippet": snippet
        })

    # If nothing survived the thresholding (edge-case), fallback to top_k by similarity
    if not sources_out:
        ranked = sorted(best_by_idx.items(), key=lambda x: x[1], reverse=True)[:top_k]
        selected_idxs = [idx for idx, _ in ranked]
        retrieved = [chunks[i] for i in selected_idxs]
        labeled_context = "\n\n".join(f"[{i+1}] {txt}" for i, txt in enumerate(retrieved))
        allowed_labels = [str(i + 1) for i in range(len(retrieved))]
        labels_hint = ", ".join(f"[{x}]" for x in allowed_labels)
        sources_out = []
        for n, idx in enumerate(selected_idxs, start=1):
            if meta and 0 <= idx < len(meta):
                m = meta[idx]
                file = m.get("file")
                p_start = m.get("page_start")
                p_end = m.get("page_end")
                snippet = m.get("snippet") or (retrieved[n - 1][:240] + "…"
                                               if len(retrieved[n - 1]) > 240 else retrieved[n - 1])
            else:
                file = None
                p_start = None
                p_end = None
                snippet = retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1]
            sources_out.append({
                "label": n,
                "file": file,
                "page_start": p_start,
                "page_end": p_end,
                "snippet": snippet
            })

    # -------- Pass 1: QUOTES (extractive) --------
    quotes_system = f"""
Είσαι συστηματικός επιμελητής πηγών. Δούλεψε ΜΟΝΟ από τις ΠΗΓΕΣ.
Στόχος: Εξήγαγε 2–4 σύντομα αποσπάσματα (<= 30 λέξεις) που απαντούν άμεσα την ΕΡΩΤΗΣΗ.

Κανόνες:
- ΜΟΝΟ άμεσες παραθέσεις (quotes) από τις ΠΗΓΕΣ. Καμία παράφραση εδώ.
- Κάθε quote να τελειώνει με ΣΩΣΤΗ ετικέτα από τα επιτρεπόμενα: {labels_hint}
- Κράτα μόνο τα πιο σχετικά σημεία, χωρίς επαναλήψεις.

Μορφή (αυστηρά):
- "…απόσπασμα…" [n]
- "…απόσπασμα…" [n]
- (το πολύ 4 γραμμές)
"""
    quotes_prompt = f"""{quotes_system}

ΕΡΩΤΗΣΗ:
{question}

ΠΗΓΕΣ:
{labeled_context}

ΑΠΟΣΠΑΣΜΑΤΑ:"""
    quotes = query_model(quotes_prompt)

    # -------- Pass 2: FINAL ANSWER (compressive from quotes) --------
    if EXTRACTIVE_ONLY:
        # quotes-only mode (no paraphrase)
        lines = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
        final_answer = "Συνοπτικά:\n" + "\n".join(lines[:4]) if lines else "Δεν ξέρω."
    else:
        answer_system = f"""
Συντάκτης απαντήσεων, αυστηρά βασισμένων στα αποσπάσματα.
Σύνθεση ΜΟΝΟ από τα ΠΑΡΑΠΑΝΩ QUOTES, χωρίς νέα γνώση.

Κανόνες:
- Μήκος: έως 2 σύντομες προτάσεις συνολικά. Προαιρετικά έως 3 κουκκίδες.
- Κάθε πρόταση/κουκκίδα με ισχυρισμό τελειώνει με έγκυρη ετικέτα από: {labels_hint}.
- Αν δεν μπορείς να τεκμηριώσεις, απάντησε: "Δεν ξέρω."
"""
        final_prompt = f"""{answer_system}

ΕΡΩΤΗΣΗ:
{question}

QUOTES:
{quotes}

ΤΕΛΙΚΗ ΑΠΑΝΤΗΣΗ:"""
        final_answer = query_model(final_prompt).strip()

    # Safety net: if model forgot citations, attach labels seen in QUOTES
    final_answer = ensure_citations_from_quotes(final_answer, quotes, allowed_labels)

    # Strict fallback: if still no [n], refuse
    if STRICT_MODE and not re.search(r"\[\d+\]", final_answer):
        final_answer = "Δεν ξέρω."

    # --- Keep only sources that were actually cited in the final answer ---
    used_labels = sorted(set(int(n) for n in re.findall(r"\[(\d+)\]", final_answer)))
    if used_labels:
        filtered_sources = [s for s in sources_out if s["label"] in used_labels]
        order = {lab: i for i, lab in enumerate(used_labels)}
        filtered_sources.sort(key=lambda s: order.get(s["label"], 9999))
    else:
        filtered_sources = []

    return final_answer, filtered_sources


# -----------------------------
# File serving (PDFs only)
# -----------------------------
@app.route("/files/<domain>/<path:filename>", methods=["GET"])
def serve_file(domain, filename):
    # only allow .pdf
    if not filename.lower().endswith(".pdf"):
        abort(403)
    domain_path = os.path.join(BASE_FOLDER, domain)
    if not os.path.isdir(domain_path):
        abort(404)
    safe_path = safe_join(domain_path, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)
    return send_from_directory(domain_path, filename, as_attachment=False)


# -----------------------------
# Flask endpoints
# -----------------------------
def require_api_key():
    key = request.headers.get("x-api-key")
    if not API_KEY or API_KEY == "changeme":
        print("⚠️ WARNING: API_KEY is default. Set API_KEY env var before public exposure!")
    elif key != API_KEY:
        abort(401, description="Unauthorized")


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "model": MODEL_NAME, "domains": list(domain_data.keys())})


@app.route("/ask", methods=["POST"])
def ask():
    require_api_key()
    data = request.get_json(force=True)
    question = data.get("question", "")
    domain_in = data.get("domain", "default")
    # Case-insensitive domain mapping to existing loaded domains
    matched_domain = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), domain_in)
    print(f"Question for domain [{matched_domain}]: {question}")
    answer, sources = answer_question(matched_domain, question)
    return jsonify({"answer": answer, "sources": sources})


# -----------------------------
# DOWNLOAD PDF PAGES (excerpt)
# -----------------------------
@app.route("/files/clip", methods=["GET"])
def serve_pdf_clip():
    domain = request.args.get("domain")
    filename = request.args.get("file")
    page_from = request.args.get("from", type=int)
    page_to = request.args.get("to", type=int)

    if not (domain and filename and page_from):
        abort(400, description="Missing domain/file/from")

    if not filename.lower().endswith(".pdf"):
        abort(403)

    domain_path = os.path.join(BASE_FOLDER, domain)
    if not os.path.isdir(domain_path):
        abort(404)

    safe_path = safe_join(domain_path, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)

    if page_to is None:
        page_to = page_from

    # PyPDF2 uses 0-based internally; our meta is 1-based for users
    try:
        reader = PyPDF2.PdfReader(safe_path)
        total = len(reader.pages)
        start = max(1, page_from)
        end = max(start, page_to)
        start0 = min(total, start) - 1
        end0 = min(total, end) - 1

        writer = PyPDF2.PdfWriter()
        for p in range(start0, end0 + 1):
            writer.add_page(reader.pages[p])

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
    except Exception as e:
        print("Clip error:", e)
        abort(500, description="Failed to clip PDF")

    out_name = f"{os.path.splitext(filename)[0]}_p{start}-p{end}.pdf"
    return send_file(buf, mimetype="application/pdf", as_attachment=True, download_name=out_name)


def load_domains():
    if not os.path.isdir(BASE_FOLDER):
        print(f"⚠️ BASE_FOLDER not found: {BASE_FOLDER}")
        return
    found_any = False
    for domain in os.listdir(BASE_FOLDER):
        folder = os.path.join(BASE_FOLDER, domain)
        if os.path.isdir(folder):
            build_domain_index(domain, folder)
            found_any = True
    if not found_any:
        print(f"⚠️ No domain subfolders found in {BASE_FOLDER}. Create e.g.:")
        print(os.path.join(BASE_FOLDER, "cybersecurity"))
        print(os.path.join(BASE_FOLDER, "digitalMarketing"))


load_domains()

if __name__ == "__main__":
    # 0.0.0.0 so ngrok can reach it
    app.run(host="0.0.0.0", port=5000)
