# ai-testpy_VLLM_TEACHER_ANS.py
# Stable server with:
# - legacy flat domains supported
# - nested Courses/<Course>/YEAR_X/SEM_Y/<Lesson> supported
# - year-scope escalation via scope="year"
# - ✅ accepts year alias: Courses__<Course>__YEAR_X
# - safe file serving for nested domain folders

from flask import Flask, request, jsonify, abort, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import safe_join
import os, io, re, time, json, hashlib, shutil, unicodedata
import numpy as np
import requests
import PyPDF2

try:
    import faiss
except Exception:
    faiss = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import torch
except Exception:
    torch = None


app = Flask(__name__)
CORS(app)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# === Core Config (keep compatible) ===
DEFAULT_BASE = os.path.join(SCRIPT_DIR, "domains")
BASE_FOLDER = os.getenv("BASE_FOLDER", DEFAULT_BASE)
os.makedirs(BASE_FOLDER, exist_ok=True)

API_KEY = os.getenv("API_KEY", "").strip()
if not API_KEY:
    print("WARNING: API_KEY is empty. Set it via env for production use.")

VLLM_URL = os.getenv("VLLM_URL", "http://127.0.0.1:11434/api/chat")
VLLM_MODEL = os.getenv("VLLM_MODEL", "llama3.1")
VLLM_TIMEOUT = int(os.getenv("VLLM_TIMEOUT", "120"))

MODEL_NAME = os.getenv("MODEL_NAME", "vLLM-chat")

TEMPERATURE = float(os.getenv("TEMPERATURE", "0.2"))
TOP_P = float(os.getenv("TOP_P", "0.9"))
MAX_TOKENS = int(os.getenv("MAX_TOKENS", "240"))

# vLLM token budgeting (prevents 400s when prompt is large)
VLLM_MAX_CONTEXT = int(os.getenv("VLLM_MAX_CONTEXT", "2048"))
VLLM_SAFETY_MARGIN = int(os.getenv("VLLM_SAFETY_MARGIN", "256"))
VLLM_MIN_COMPLETION = int(os.getenv("VLLM_MIN_COMPLETION", "180"))
VLLM_MAX_COMPLETION_CAP = int(os.getenv("VLLM_MAX_COMPLETION_CAP", str(MAX_TOKENS)))

# Retrieval tuning
CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "280"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "20"))
R_TOP_K = int(os.getenv("R_TOP_K", "4"))
NUM_ALTERNATE_QUERIES = int(os.getenv("NUM_ALTERNATE_QUERIES", "1"))
MIN_SIM_THRESHOLD = float(os.getenv("MIN_SIM_THRESHOLD", "0.30"))

STRICT_MODE = os.getenv("STRICT_MODE", "true").lower() in ("1", "true", "yes")

# ✅ Strict fallback phrase (must match exactly)
STRICT_NO_EVIDENCE_PHRASE = "Δεν υπάρχουν αρκετά στοιχεία στις πηγές."
EXTRACTIVE_ONLY = os.getenv("EXTRACTIVE_ONLY", "false").lower() in ("1", "true", "yes")

EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "64"))
MIN_WORDS_TO_INDEX = int(os.getenv("MIN_WORDS_TO_INDEX", "20"))

# Teacher answers (per-domain PDF)
TEACHER_MIN_SIM = float(os.getenv("TEACHER_MIN_SIM", "0.55"))

# { domain_id: { "questions": [...], "answers": [...],
#               "embeddings": np.ndarray, "pages": [...], "file": "..." } }
teacher_answers_cache = {}

# Index cache
INDEX_CACHE_DIR = os.getenv("INDEX_CACHE_DIR", os.path.join(SCRIPT_DIR, ".index_cache"))
os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

CACHE_ADMIN_KEY = os.getenv("CACHE_ADMIN_KEY", "").strip()

SUPPORTED_EXTS = (".pdf", ".txt", ".docx", ".pptx")

# Embedding model
EMBEDDING_DEVICE = os.getenv("EMBEDDING_DEVICE", "").strip().lower()
if not EMBEDDING_DEVICE:
    if torch is not None and torch.cuda.is_available():
        EMBEDDING_DEVICE = "cpu"
    else:
        EMBEDDING_DEVICE = "cpu"

if SentenceTransformer is None:
    raise RuntimeError("sentence-transformers is required")
if faiss is None:
    raise RuntimeError("faiss is required")

embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2", device=EMBEDDING_DEVICE)

try:
    embedding_model.encode(["warmup"], convert_to_numpy=True, normalize_embeddings=True)
    print("✅ Embedding model warmup complete.")
except Exception as e:
    print("⚠️ Embedding warmup failed: ", e)

PROMPT_LANG = (
    "Απάντησε πάντα στα Ελληνικά, με φυσική και κατανοητή γλώσσα. "
    "Αν η ερώτηση είναι σε άλλη γλώσσα, μετέφρασέ την συνοπτικά και απάντησε στα Ελληνικά."
)

# =========================
# In-memory indexes
# =========================
# domain_data:
# { domainId: { "chunks": [...], "index": faiss.IndexFlatIP, "meta": [...],
#              "embeddings": np.ndarray, "folder": "...", "kind": "lesson|year",
#              "children": [...]} }
domain_data = {}

# domainId -> folder path
domain_folders = {}

# mapping for nested Courses/<Course>/YEAR_X/SEM_Y/<Lesson>
lesson_to_year = {}
year_alias_to_year = {}


# =========================
# Helpers
# =========================
def require_api_key():
    key = request.headers.get("x-api-key") or ""
    if not API_KEY:
        # For local dev, we allow missing api-key but log it
        if key:
            print("⚠️ API_KEY not configured on server but client sent a key.")
        return
    if key != API_KEY:
        abort(401, description="Unauthorized")


def normalize_text(s: str) -> str:
    if not s:
        return ""
    s = unicodedata.normalize("NFKC", s)
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def rough_tokenize(text: str) -> int:
    # crude: count words as tokens
    return len((text or "").split())


def sliding_word_chunks(words, max_tokens, overlap_tokens):
    chunks = []
    cur = []
    cur_len = 0
    for w in words:
        cur.append(w)
        cur_len += 1
        if cur_len >= max_tokens:
            chunks.append(" ".join(cur))
            if overlap_tokens > 0:
                cur = cur[-overlap_tokens:]
                cur_len = len(cur)
            else:
                cur = []
                cur_len = 0
    if cur:
        chunks.append(" ".join(cur))
    return chunks


def safe_read_pdf_text(path):
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            texts = []
            for page in reader.pages:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t:
                    texts.append(t)
            return "\n".join(texts)
    except Exception as e:
        print(f"Error reading PDF {path}: {e}")
        return ""


def safe_read_docx_text(path):
    if docx is None:
        return ""
    try:
        d = docx.Document(path)
        return "\n".join([p.text for p in d.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX {path}: {e}")
        return ""


def safe_read_pptx_text(path):
    if Presentation is None:
        return ""
    try:
        pres = Presentation(path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"Error reading PPTX {path}: {e}")
        return ""


def read_file_text(path):
    lower = path.lower()
    if lower.endswith(".pdf"):
        return safe_read_pdf_text(path)
    elif lower.endswith(".docx"):
        return safe_read_docx_text(path)
    elif lower.endswith(".pptx"):
        return safe_read_pptx_text(path)
    else:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file {path}: {e}")
            return ""


def embed_texts(texts):
    if not texts:
        return np.zeros((0, 384), dtype=np.float32)
    all_vecs = []
    for i in range(0, len(texts), EMBED_BATCH_SIZE):
        batch = texts[i : i + EMBED_BATCH_SIZE]
        emb = embedding_model.encode(batch, convert_to_numpy=True, normalize_embeddings=True)
        all_vecs.append(emb.astype(np.float32))
    return np.vstack(all_vecs)


def load_cached_file(domain_id: str, fpath: str, chunk_tokens: int, overlap_tokens: int):
    """
    Load cached chunks+embeddings for a single file, if present.
    """
    rel = os.path.relpath(fpath, BASE_FOLDER)
    key_str = f"{domain_id}::{rel}::{chunk_tokens}::{overlap_tokens}"
    h = hashlib.sha256(key_str.encode("utf-8", errors="ignore")).hexdigest()
    folder = os.path.join(INDEX_CACHE_DIR, h)
    meta_path = os.path.join(folder, "meta.json")
    emb_path = os.path.join(folder, "embeddings.npy")
    if not (os.path.isfile(meta_path) and os.path.isfile(emb_path)):
        return None
    try:
        with open(meta_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        embeddings = np.load(emb_path)
    except Exception as e:
        print(f"Failed to load cached file index for {fpath}: {e}")
        return None
    return data, embeddings


def save_cached_file(domain_id: str, fpath: str, chunk_tokens: int, overlap_tokens: int, data, embeddings):
    """
    Save chunks+embeddings for a single file.
    """
    rel = os.path.relpath(fpath, BASE_FOLDER)
    key_str = f"{domain_id}::{rel}::{chunk_tokens}::{overlap_tokens}"
    h = hashlib.sha256(key_str.encode("utf-8", errors="ignore")).hexdigest()
    folder = os.path.join(INDEX_CACHE_DIR, h)
    os.makedirs(folder, exist_ok=True)
    meta_path = os.path.join(folder, "meta.json")
    emb_path = os.path.join(folder, "embeddings.npy")
    try:
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False)
        np.save(emb_path, embeddings)
    except Exception as e:
        print(f"Failed to save cached file index for {fpath}: {e}")


def parse_file_to_chunks(fname: str, folder: str):
    """
    Parse a single file into chunks and metadata.
    """
    path = os.path.join(folder, fname)
    text = read_file_text(path)
    text = normalize_text(text)
    if not text:
        return None, None, None

    words = text.split()
    if len(words) < MIN_WORDS_TO_INDEX:
        return None, None, None

    w_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
    chunks = []
    metas = []

    for ch in w_chunks:
        chunks.append(ch)
        metas.append(
            {
                "file": fname,
                "rel": os.path.relpath(path, folder),
                "path": path,
                "page": None,
            }
        )

    return path, chunks, metas


def build_domain_index(domain: str, folder: str):
    t0 = time.time()

    files = [f for f in os.listdir(folder) if f.lower().endswith(SUPPORTED_EXTS)]
    files.sort()
    if not files:
        print(f"⚠️ No usable files for '{domain}' in {folder}")
        return

    all_chunks = []
    all_metas = []
    all_embeddings_parts = []

    misses = []
    for f in files:
        path = os.path.join(folder, f)
        cached = load_cached_file(domain, path, CHUNK_TOKENS, CHUNK_OVERLAP)

        if cached is not None:
            data, emb = cached
            chunks = data.get("chunks") or []
            metas = data.get("meta") or []

            if chunks and metas and emb is not None and len(chunks) == emb.shape[0]:
                print(f"  ↪ reused cache for {f}: {len(chunks)} chunks")
                all_chunks.extend(chunks)
                all_metas.extend(metas)
                all_embeddings_parts.append(emb.astype(np.float32))
            else:
                # cache is corrupt or empty, re-index this file
                misses.append(f)
        else:
            misses.append(f)

    parsed = []
    if misses:
        print(f"  ↪ indexing {len(misses)} files (no cache)")
        for f in misses:
            result = parse_file_to_chunks(f, folder)
            if not result:
                continue
            path, chunks, metas = result
            if chunks and metas:
                parsed.append((path, chunks, metas))

        if parsed:
            flat_chunks = []
            boundaries = []
            cursor = 0

            for path, chunks, metas in parsed:
                start = cursor
                flat_chunks.extend(chunks)
                cursor += len(chunks)
                boundaries.append((path, start, cursor, metas))

            emb = embed_texts(flat_chunks)

            for (path, start, end, metas) in boundaries:
                file_chunks = flat_chunks[start:end]
                file_emb = emb[start:end]

                # save per-file cache
                save_cached_file(
                    domain,
                    path,
                    CHUNK_TOKENS,
                    CHUNK_OVERLAP,
                    {"chunks": file_chunks, "meta": metas},
                    file_emb,
                )

                all_chunks.extend(file_chunks)
                all_metas.extend(metas)
                all_embeddings_parts.append(file_emb.astype(np.float32))

    if all_embeddings_parts:
        embeddings = np.vstack(all_embeddings_parts).astype(np.float32)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)
    else:
        embeddings = np.zeros((0, 384), dtype=np.float32)
        index = faiss.IndexFlatIP(384)

    domain_data[domain] = {
        "chunks": all_chunks,
        "meta": all_metas,
        "embeddings": embeddings,
        "index": index,
        "folder": folder,
        "kind": "lesson",
        "children": [],
    }
    domain_folders[domain] = folder

    dt = time.time() - t0
    print(f"✅ Indexed {len(all_chunks)} chunks for '{domain}' in {dt:.2f}s.")



def discover_domain_folders():
    """
    Discover both:
    - flat domains: BASE_FOLDER/<domainId>
    - nested courses: BASE_FOLDER/Courses/<Course>/YEAR_X/SEM_Y/<Lesson>
    """
    discovered = []

    # 1) flat domains
    for name in os.listdir(BASE_FOLDER):
        if name == "Courses":
            continue
        path = os.path.join(BASE_FOLDER, name)
        if os.path.isdir(path):
            discovered.append((name, path, name))

    # 2) nested courses
    courses_root = os.path.join(BASE_FOLDER, "Courses")
    if os.path.isdir(courses_root):
        for course_name in os.listdir(courses_root):
            course_path = os.path.join(courses_root, course_name)
            if not os.path.isdir(course_path):
                continue

            for year_name in os.listdir(course_path):
                if not re.match(r"YEAR_\d+", year_name, re.IGNORECASE):
                    continue
                year_path = os.path.join(course_path, year_name)
                if not os.path.isdir(year_path):
                    continue

                for sem_name in os.listdir(year_path):
                    sem_path = os.path.join(year_path, sem_name)
                    if not os.path.isdir(sem_path):
                        continue

                    for lesson_name in os.listdir(sem_path):
                        lesson_path = os.path.join(sem_path, lesson_name)
                        if not os.path.isdir(lesson_path):
                            continue

                        domain_id = f"Courses__{course_name}__{year_name}__{sem_name}__{lesson_name}"
                        discovered.append((domain_id, lesson_path, os.path.relpath(lesson_path, BASE_FOLDER)))

    return discovered


def build_year_aggregates():
    """
    Build YEAR__ indices from existing lesson domains.
    """
    lesson_to_year.clear()
    year_alias_to_year.clear()

    # group lessons by (course, year)
    grouped = {}  # (course, year) -> [domain_id]
    for domain_id in list(domain_data.keys()):
        if not domain_id.startswith("Courses__"):
            continue
        parts = domain_id.split("__")
        if len(parts) < 4:
            continue
        _, course_name, year_name, sem_name, *rest = parts
        key = (course_name, year_name)
        grouped.setdefault(key, []).append(domain_id)

    for (course_name, year_name), children in grouped.items():
        course_slug = course_name.replace(" ", "_")
        year_domain = f"YEAR__{course_slug}__{year_name}"

        # map lessons -> year
        for lesson_domain in children:
            lesson_to_year[lesson_domain] = year_domain

        # alias: Courses__Course__YEAR_X -> YEAR__Course__YEAR_X
        alias_key = f"Courses__{course_name}__{year_name}".replace(" ", "_")
        year_alias_to_year[alias_key] = year_domain

        # build aggregated index
        chunks_all = []
        meta_all = []
        embeddings_all = []

        for child in children:
            info = domain_data.get(child)
            if not info:
                continue
            if info["embeddings"].shape[0] == 0:
                continue
            childs_chunks = info["chunks"]
            childs_meta = info["meta"]
            childs_emb = info["embeddings"]
            chunks_all.extend(childs_chunks)
            meta_all.extend(childs_meta)
            embeddings_all.append(childs_emb)

        if not embeddings_all:
            continue

        embeddings = np.vstack(embeddings_all)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)

        domain_data[year_domain] = {
            "chunks": chunks_all,
            "meta": meta_all,
            "embeddings": embeddings,
            "index": index,
            "folder": None,
            "kind": "year",
            "children": children,
        }

        print(f"🧩 YEAR index '{year_domain}' built from {len(children)} lesson domains.")


def load_domains():
    domain_data.clear()
    domain_folders.clear()
    lesson_to_year.clear()
    year_alias_to_year.clear()

    discovered = discover_domain_folders()

    for domain_id, folder, rel in discovered:
        build_domain_index(domain_id, folder)

    build_year_aggregates()


# =========================
# Retrieval
# =========================
def retrieve_chunks(domain_id: str, query: str, top_k: int = R_TOP_K, alt_queries=None):
    """
    Retrieve top_k chunks for a query (and optional alternate queries) from domain.
    Returns: list of (chunk_text, score, meta)
    """
    info = domain_data.get(domain_id)
    if not info or "index" not in info:
        raise ValueError(f"Unknown or unindexed domain: {domain_id}")

    index = info["index"]
    embeddings = info["embeddings"]
    chunks = info["chunks"]
    meta = info["meta"]

    if embeddings.shape[0] == 0:
        return []

    queries = [query] + (alt_queries or [])
    q_emb = embed_texts(queries)
    D, I = index.search(q_emb, top_k)

    results = []
    used = set()
    for row in range(D.shape[0]):
        for col in range(D.shape[1]):
            idx = int(I[row, col])
            if idx < 0 or idx >= len(chunks):
                continue
            if idx in used:
                continue
            score = float(D[row, col])
            if score < MIN_SIM_THRESHOLD:
                continue
            used.add(idx)
            results.append((chunks[idx], score, meta[idx]))

    results.sort(key=lambda x: x[1], reverse=True)
    return results[:top_k]


def auto_extract_quotes(chunks):
    """
    Extract short bullet-like segments from chunks for quote-based answers
    when vLLM is not available.
    """
    lines = []
    for ch, _, _ in chunks:
        for ln in re.split(r"[.?!]\s+", ch):
            ln = ln.strip()
            if not ln:
                continue
            if len(ln.split()) < 200:
                lines.append("- " + ln.strip())
    return "\n".join(lines)


def build_prompt_from_chunks(question, retrieved, chat_messages=None):
    """
    Build the final prompt given retrieved chunks and a question.
    """
    context_parts = []
    sources = []

    for i, (chunk, score, m) in enumerate(retrieved, start=1):
        file = m.get("file")
        page = m.get("page")
        snippet = chunk
        if len(snippet) > 240:
            snippet = snippet[:240] + "…"

        sources.append(
            {
                "label": i,
                "file": file,
                "page_start": page,
                "page_end": page,
                "snippet": snippet,
            }
        )

        context_parts.append(f"[{i}] ({file}, σελίδα {page}) {chunk}")

    history_str = ""
    if chat_messages:
        history_items = []
        for msg in chat_messages[-4:]:
            role = msg.get("role")
            content = msg.get("content") or ""
            if role == "user":
                history_items.append(f"Χρήστης: {content}")
            elif role == "assistant":
                history_items.append(f"Βοηθός: {content}")
        if history_items:
            history_str = "\n\nΙστορικό συνομιλίας:\n" + "\n".join(history_items)

    context_text = "\n\n".join(context_parts) if context_parts else ""
    # ✅ Strict, anti-meta instructions (the system prompt reinforces this too)
    instructions = (
        f"{PROMPT_LANG}\n"
        "Απάντησε ΜΟΝΟ στα Ελληνικά.\n"
        "Απάντησε κατευθείαν, χωρίς πρόλογο/ευχαριστίες/μετα-σχόλια.\n"
        "Χρησιμοποίησε ΜΟΝΟ τις παρακάτω πληροφορίες από τα έγγραφα.\n"
        f"Αν οι πληροφορίες δεν επαρκούν, απάντησε ΑΚΡΙΒΩΣ: \"{STRICT_NO_EVIDENCE_PHRASE}\"\n"
        "Μην επαναλαμβάνεις τις οδηγίες.\n"
    )

    full_prompt = (
        instructions
        + history_str
        + ("\n\nΠληροφορίες από έγγραφα:\n" + context_text if context_text else "")
        + f"\n\nΕρώτηση: {question}\n\nΑπάντηση:"
    )

    return full_prompt, sources


# =========================
# LLM
# =========================
def _approx_token_count(text: str) -> int:
    if not text:
        return 0
    return max(1, int(len(text) / 4))


def _truncate_text_to_token_budget(text: str, max_tokens: int) -> str:
    """\
    Truncate while preserving BOTH the beginning (instructions) and the end (question).

    The old behavior was `text[:n]` which could accidentally cut off the user's question
    (usually placed at the end of the prompt), causing the model to reply with
    "Παρακαλώ, δώσε μου την ερώτηση...".
    """
    if not text:
        return text
    approx_tokens = _approx_token_count(text)
    if approx_tokens <= max_tokens:
        return text

    ratio = max_tokens / float(approx_tokens)
    n_chars = max(0, int(len(text) * ratio))
    if n_chars <= 0:
        return ""

    # Keep a head + tail slice to preserve instructions + the actual question.
    head_chars = min(1200, max(0, n_chars // 3))
    tail_chars = max(0, n_chars - head_chars)
    if tail_chars <= 0:
        return text[:n_chars]
    return text[:head_chars] + "\n\n...\n\n" + text[-tail_chars:]


def query_model(prompt: str) -> str:
    """
    Call vLLM chat endpoint with token budgeting.
    """
    est_prompt_tokens = _approx_token_count(prompt)
    max_context = VLLM_MAX_CONTEXT
    safety = VLLM_SAFETY_MARGIN

    avail_for_completion = max_context - est_prompt_tokens - safety
    avail_for_completion = max(avail_for_completion, VLLM_MIN_COMPLETION)
    avail_for_completion = min(avail_for_completion, VLLM_MAX_COMPLETION_CAP)

    # Make sure prompt is within context budget
    prompt_trimmed = _truncate_text_to_token_budget(prompt, max_context - safety - avail_for_completion)

    payload = {
        "model": VLLM_MODEL,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Απάντα ΜΟΝΟ στα Ελληνικά. "
                    "Απάντα κατευθείαν στην ερώτηση, χωρίς πρόλογο/ευχαριστίες/μετα-σχόλια. "
                    "Χρησιμοποίησε ΜΟΝΟ τις πληροφορίες που σου δίνονται από τα αποσπάσματα. "
                    f"Αν δεν υπάρχουν αρκετά στοιχεία στα αποσπάσματα, απάντησε ΑΚΡΙΒΩΣ: \"{STRICT_NO_EVIDENCE_PHRASE}\". "
                    "Μην επαναλαμβάνεις οδηγίες."
                ),
            },
            {"role": "user", "content": prompt_trimmed},
        ],
        "stream": False,
        "options": {
            "temperature": TEMPERATURE,
            "top_p": TOP_P,
            "num_predict": avail_for_completion,
        },
    }

    resp = requests.post(VLLM_URL, json=payload, timeout=VLLM_TIMEOUT)
    if resp.status_code != 200:
        raise RuntimeError(f"vLLM HTTP {resp.status_code}: {resp.text[:500]}")
    data = resp.json()
    content = data["message"]["content"]
    return content.strip() or "Δεν μπόρεσα να παραγάγω απάντηση."


# =========================
# Smart relevance guard (lexical sanity check)
# =========================
import re
import unicodedata


_GREEK_STOPWORDS = {
    "τι", "ποιος", "ποια", "ποιο", "ποιες", "ποιοι", "πώς", "πως", "γιατί", "γιατι",
    "είναι", "ειναι", "ή", "η", "ο", "το", "τα", "τους", "τις", "των", "στο", "στη", "στην",
    "σε", "με", "και", "να", "θα", "που", "από", "απο", "πως", "πάνω", "κάτω", "μια", "μία",
    "ένα", "ενα", "ένας", "ενας", "μιας", "μια", "μου", "σου", "του", "της", "μας", "σας",
}


def _strip_accents(s: str) -> str:
    if not s:
        return ""
    s = unicodedata.normalize("NFD", s)
    s = "".join(ch for ch in s if unicodedata.category(ch) != "Mn")
    return unicodedata.normalize("NFC", s)


def _tokenize_keywords(text: str) -> list[str]:
    """Tokenize Greek-ish text, remove accents, lowercase, drop stopwords."""
    if not text:
        return []
    t = _strip_accents(text.lower())
    tokens = re.findall(r"[\wάέήίόύώΑ-Ωα-ω]+", t, flags=re.UNICODE)
    out = []
    for tok in tokens:
        tok = tok.strip("_")
        if len(tok) < 4:
            continue
        if tok in _GREEK_STOPWORDS:
            continue
        out.append(tok)
    # de-dup while preserving order
    seen = set()
    dedup = []
    for t in out:
        if t not in seen:
            seen.add(t)
            dedup.append(t)
    return dedup


def _lexical_relevance_ok(question: str, retrieved: list, min_hits: int = 1) -> bool:
    """Return True if retrieved chunks contain at least `min_hits` keyword hits from question."""
    kws = _tokenize_keywords(question)
    if not kws:
        # If we can't extract keywords, don't block.
        return True

    # Build a normalized haystack of top retrieved chunks
    top_text = "\n".join([(c or "") for (c, _, _) in retrieved[: max(6, len(retrieved))]])
    hay = _strip_accents(top_text.lower())

    hits = 0
    for kw in kws[:8]:  # check up to 8 keywords
        # Be tolerant to PDF extraction artifacts (line breaks/hyphenation) by also
        # allowing a short prefix match.
        if kw in hay or (len(kw) >= 7 and kw[:6] in hay):
            hits += 1
            if hits >= min_hits:
                return True
    return False


def answer_question(domain_id: str, question: str, chat_messages=None):
    """
    Main QA pipeline for a given domain.
    Keeps original behavior, including quote-based fallback when vLLM fails.
    """
    alt_queries = []
    if NUM_ALTERNATE_QUERIES > 0:
        alt_queries = [question.upper()] if NUM_ALTERNATE_QUERIES >= 1 else []

    retrieved = retrieve_chunks(domain_id, question, top_k=R_TOP_K, alt_queries=alt_queries)
    if not retrieved:
        if STRICT_MODE:
            return (STRICT_NO_EVIDENCE_PHRASE, [])
        else:
            prompt = (
                f"{PROMPT_LANG}\n"
                "Δυστυχώς δεν βρήκα σχετικές πληροφορίες στα έγγραφα, αλλά προσπάθησε να απαντήσεις "
                "με βάση τις γενικές σου γνώσεις.\n\n"
                f"Ερώτηση: {question}\n\nΑπάντηση:"
            )
            try:
                ans = query_model(prompt)
            except Exception as e:
                print("⚠️ vLLM call failed in no-context path:", e)
                ans = "Δεν μπόρεσα να απαντήσω."
            return ans, []

    # ✅ Lexical sanity check: if retrieval seems unrelated, prefer strict fallback.
    # This prevents answers that are "about something else" because embeddings matched noise.
    if STRICT_MODE and not _lexical_relevance_ok(question, retrieved, min_hits=1):
        return (STRICT_NO_EVIDENCE_PHRASE, [])

    prompt, sources = build_prompt_from_chunks(question, retrieved, chat_messages=chat_messages)

    if EXTRACTIVE_ONLY:
        joined = "\n\n".join([c for (c, _, _) in retrieved])
        return joined, sources

    # Use quote-based fallback if vLLM fails
    quotes = auto_extract_quotes(retrieved)
    full_prompt = prompt
    try:
        final_answer = query_model(full_prompt)
    except Exception as e:
        print("⚠️ vLLM call failed; returning extractive answer:", e)
        lines_ = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
        final_answer = "Συνοπτικά:\n" + "\n".join(lines_[:4]) if lines_ else "Δεν ξέρω."

    # Optionally, you might want to filter sources here (unchanged from original)
    filtered_sources = sources

    return final_answer, filtered_sources


def _require_cache_admin():
    if not CACHE_ADMIN_KEY:
        abort(403, description="Cache admin key not configured")
    key = request.headers.get("x-cache-admin") or ""
    if key != CACHE_ADMIN_KEY:
        abort(403, description="Invalid cache admin key")


# =========================
# Domain folder resolving
# =========================
def _resolve_domain_folder(domain: str):
    if domain in domain_folders:
        return domain_folders[domain]
    d2 = next((k for k in domain_folders.keys() if k.lower() == (domain or "").lower()), None)
    if d2:
        return domain_folders[d2]
    return None


def load_teacher_answers_for_domain(domain: str):
    """
    Load and cache teacher answers for a given domain.

    Expected file name inside the domain folder:
      {domain}_TEACHER_ANSWERS.pdf

    Example:
      domain = "Courses__Psychology__YEAR_2__SEM_3__ΨΥΧΓΕ3"
      file   = "Courses__Psychology__YEAR_2__SEM_3__ΨΥΧΓΕ3_TEACHER_ANSWERS.pdf"
    """
    if not domain:
        return None

    # Cache hit
    if domain in teacher_answers_cache:
        return teacher_answers_cache[domain]

    folder = _resolve_domain_folder(domain)
    if not folder:
        teacher_answers_cache[domain] = None
        return None

    teacher_fname = f"{domain}_TEACHER_ANSWERS.pdf"
    teacher_path = os.path.join(folder, teacher_fname)
    if not os.path.isfile(teacher_path):
        teacher_answers_cache[domain] = None
        return None

    questions = []
    answers = []
    pages = []

    try:
        reader = PyPDF2.PdfReader(teacher_path)
    except Exception as e:
        print(f"⚠️ Could not read teacher PDF for domain {domain}: {e}")
        teacher_answers_cache[domain] = None
        return None

    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""

        if not text.strip():
            continue

        # Adjust these keys if your real PDF uses slightly different labels
        q_key = "Ερώτηση:"
        a_key = "Απάντηση καθηγητή:"
        d_key = "DocId:"

        q_idx = text.find(q_key)
        a_idx = text.find(a_key)

        if q_idx == -1 or a_idx == -1 or a_idx <= q_idx:
            continue

        q_text = text[q_idx + len(q_key):a_idx].strip()

        d_idx = text.find(d_key, a_idx)
        if d_idx != -1:
            a_text = text[a_idx + len(a_key):d_idx].strip()
        else:
            a_text = text[a_idx + len(a_key):].strip()

        if q_text and a_text:
            questions.append(q_text)
            answers.append(a_text)
            pages.append(i + 1)  # 1-based page index

    if not questions:
        teacher_answers_cache[domain] = None
        return None

    try:
        emb = embedding_model.encode(
            questions,
            convert_to_numpy=True,
            normalize_embeddings=True
        ).astype(np.float32)
    except Exception as e:
        print(f"⚠️ Failed to embed teacher questions for domain {domain}: {e}")
        teacher_answers_cache[domain] = None
        return None

    teacher_answers_cache[domain] = {
        "questions": questions,
        "answers": answers,
        "embeddings": emb,
        "pages": pages,
        "file": teacher_fname,
    }
    print(f"📘 Loaded {len(questions)} teacher Q&A pairs for domain '{domain}'.")
    return teacher_answers_cache[domain]


def maybe_answer_from_teacher(domain: str, question: str):
    """
    Try to answer a question directly from the teacher answers PDF for this domain.
    Returns (answer_text, sources) or (None, None) if no good match.
    """
    if not question or not question.strip():
        return None, None

    info = load_teacher_answers_for_domain(domain)
    if not info:
        return None, None

    q_emb = embedding_model.encode(
        [question],
        convert_to_numpy=True,
        normalize_embeddings=True
    ).astype(np.float32)[0]

    emb_mat = info.get("embeddings")
    if emb_mat is None or not len(emb_mat):
        return None, None

    # cosine similarity (embeddings are normalized)
    sims = emb_mat @ q_emb
    best_idx = int(np.argmax(sims))
    best_sim = float(sims[best_idx])

    if best_sim < TEACHER_MIN_SIM:
        return None, None

    pdf_question = info["questions"][best_idx]
    pdf_answer = info["answers"][best_idx]

    # Include the teacher's question + answer in the response
    answer_text = (
        "Ερώτηση (μαθητή): "
        + pdf_question.strip()
        + "\n\nΑπάντηση (καθηγητή): "
        + pdf_answer.strip()
    )

    page = info["pages"][best_idx] if info.get("pages") else None
    teacher_file = info.get("file")

    snippet = pdf_answer
    if len(snippet) > 240:
        snippet = snippet[:240] + "…"

    sources = [{
        "label": 1,
        "file": teacher_file,
        "page_start": page,
        "page_end": page,
        "snippet": snippet,
    }]

    return answer_text, sources


@app.route("/files/clip", methods=["GET"])
def clip_pdf():
    domain = (request.args.get("domain") or "").strip()
    filename = (request.args.get("file") or "").strip()

    try:
        page_from = int(request.args.get("from") or "1")
        page_to = int(request.args.get("to") or str(page_from))
    except ValueError:
        abort(400, description="from/to must be integers")

    if not domain or not filename:
        abort(400, description="Missing domain or file")

    if not filename.lower().endswith(".pdf"):
        abort(403)

    folder = _resolve_domain_folder(domain)
    if not folder or not os.path.isdir(folder):
        abort(404)

    src_path = safe_join(folder, filename)
    if not src_path or not os.path.isfile(src_path):
        abort(404)

    reader = PyPDF2.PdfReader(src_path)
    total = len(reader.pages)
    if total <= 0:
        abort(404)

    # Convert to 0-based and clamp
    start = max(0, min(total - 1, page_from - 1))
    end = max(0, min(total - 1, page_to - 1))
    if end < start:
        start, end = end, start

    writer = PyPDF2.PdfWriter()
    for i in range(start, end + 1):
        writer.add_page(reader.pages[i])

    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)

    out_name = f"clip_{page_from}-{page_to}_{os.path.basename(filename)}"
    return send_file(buf, mimetype="application/pdf", as_attachment=True, download_name=out_name)

@app.route("/domains", methods=["GET"])
def list_domains():
    items = []
    for dom, folder in domain_folders.items():
        info = domain_data.get(dom) or {}
        items.append(
            {
                "domain": dom,
                "folder": folder,
                "kind": info.get("kind", "lesson"),
                "children": info.get("children", []),
            }
        )
    return jsonify(items)


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "model": MODEL_NAME})


@app.route("/ask", methods=["POST"])
def ask():
    require_api_key()
    data = request.get_json(force=True) or {}

    question = (data.get("question") or "").strip()
    domain_in = (data.get("domain") or "").strip()
    scope = (data.get("scope") or "lesson").strip().lower()  # lesson | year
    chat_messages = data.get("messages") or []

    if not question:
        return jsonify({"answer": "Δεν δόθηκε ερώτηση.", "sources": []})

    # match domain case-insensitively
    matched_domain = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), domain_in)

    # ✅ scope escalation (supports BOTH lesson id and year-alias id)
    if scope == "year":
        # 1) if frontend sends "Courses__Course__YEAR_X", map to year aggregate
        alias_key = domain_in.replace(" ", "_")
        if alias_key in year_alias_to_year and year_alias_to_year[alias_key] in domain_data:
            matched_domain = year_alias_to_year[alias_key]
        else:
            # 2) if it's a lesson id, map to its year aggregate
            yd = lesson_to_year.get(matched_domain)
            if yd and yd in domain_data:
                matched_domain = yd
            else:
                # 3) if frontend already sent YEAR__..., keep it
                yd2 = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), None)
                if yd2 and domain_data.get(yd2, {}).get("kind") == "year":
                    matched_domain = yd2

    print(f"Question for domain [{matched_domain}] (scope={scope}): {question}")

    # 1) Try teacher answers first (if {domain}_TEACHER_ANSWERS.pdf exists)
    teacher_answer, teacher_sources = maybe_answer_from_teacher(matched_domain, question)
    if teacher_answer is not None:
        return jsonify({
            "answer": teacher_answer,
            "sources": teacher_sources or [],
            "resolvedDomain": matched_domain,
            "teacherAnswer": True,  # identifier for frontend styling
        })

    # 2) Fallback to existing retrieval + vLLM + quotes logic (unchanged)
    answer, sources = answer_question(matched_domain, question, chat_messages=chat_messages)

    return jsonify({"answer": answer, "sources": sources, "resolvedDomain": matched_domain})


# =========================
# Cache admin (optional)
# =========================
@app.route("/cache/clear", methods=["POST"])
def cache_clear():
    _require_cache_admin()
    data = request.get_json(force=True) or {}

    domain = data.get("domain")
    reload_after = bool(data.get("reload", False))

    summary = {
        "domain": domain if isinstance(domain, str) and domain.strip() else None,
        "removed_files": 0,
        "removed_bytes": 0,
        "removed_domains": 0,
    }

    # Clear index cache directory
    for name in os.listdir(INDEX_CACHE_DIR):
        folder = os.path.join(INDEX_CACHE_DIR, name)
        if not os.path.isdir(folder):
            continue
        try:
            shutil.rmtree(folder)
            summary["removed_files"] += 1
        except Exception:
            pass
    os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

    if isinstance(domain, str) and domain.strip():
        # Clear specific domain from memory
        if domain in domain_data:
            domain_data.pop(domain, None)
            summary["removed_domains"] += 1

    if reload_after:
        load_domains()
        summary["reloaded"] = True
    else:
        summary["reloaded"] = False

    return jsonify(summary)


def clear_cache(domain=None):
    removed_files = 0
    removed_bytes = 0
    removed_domains = 0

    target = None
    if isinstance(domain, str) and domain.strip():
        target = domain.strip()

    # Clear from disk
    for name in os.listdir(INDEX_CACHE_DIR):
        folder = os.path.join(INDEX_CACHE_DIR, name)
        if not os.path.isdir(folder):
            continue
        try:
            shutil.rmtree(folder)
            removed_files += 1
        except Exception:
            pass
    os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

    # Clear from memory
    if target and target in domain_data:
        domain_data.pop(target, None)
        removed_domains += 1
    elif not target:
        domain_data.clear()
        removed_domains = -1  # signal "all"

    return {
        "removed_files": removed_files,
        "removed_bytes": removed_bytes,
        "removed_domains": removed_domains,
    }


# =========================
# Boot
# =========================
print("🔄 Loading domains...")
load_domains()
print(f"✅ Ready. Loaded {len(domain_data)} domains (including YEAR aggregates).")

if __name__ == "__main__":
    port = int(os.getenv("PORT", "5000"))
    app.run(host="0.0.0.0", port=port)
