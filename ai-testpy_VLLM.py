
# ai-testpy_NEW_INDEXING.py
# Faster, cached, incremental indexing for PDF/TXT/DOCX/PPTX domains.

from flask import Flask, request, jsonify, abort, send_from_directory, send_file
import os
import io
import re
import json
import hashlib
import unicodedata
import time
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone

import PyPDF2
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import requests
from flask_cors import CORS
from werkzeug.utils import safe_join  # correct import

# DOCX & PPTX
from docx import Document
from pptx import Presentation

# Optional: device detection
try:
    import torch  # type: ignore
except Exception:
    torch = None

app = Flask(__name__)
CORS(app)

# === Core Config ===
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DEFAULT_BASE = os.path.join(SCRIPT_DIR, "domains")
BASE_FOLDER = os.getenv("BASE_FOLDER", DEFAULT_BASE)
os.makedirs(BASE_FOLDER, exist_ok=True)

OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
VLLM_URL = os.getenv("VLLM_URL", "http://127.0.0.1:8000/v1/chat/completions")
VLLM_MODEL = os.getenv("VLLM_MODEL", "meta-llama/Meta-Llama-3.1-8B-Instruct")
VLLM_TIMEOUT = int(os.getenv("VLLM_TIMEOUT", "60"))
MAX_TOKENS = int(os.getenv("MAX_TOKENS", os.getenv("NUM_PREDICT", "200")))

MODEL_NAME = os.getenv("MODEL_NAME", "llama3.1")
API_KEY = os.getenv("API_KEY", "changeme")  # set a strong secret before exposing
CACHE_ADMIN_KEY = os.getenv("CACHE_ADMIN_KEY", "")  # set to enable cache admin endpoints

NUM_CTX = int(os.getenv("NUM_CTX", "4096"))
NUM_PREDICT = int(os.getenv("NUM_PREDICT", "400"))
TEMPERATURE = float(os.getenv("TEMPERATURE", "0.2"))
TOP_P = float(os.getenv("TOP_P", "0.9"))

# === Retrieval / Index Tuning ===
CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "280"))   # increased default for fewer chunks
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "20"))  # reduced default overlap
R_TOP_K = int(os.getenv("R_TOP_K", "4"))
NUM_ALTERNATE_QUERIES = int(os.getenv("NUM_ALTERNATE_QUERIES", "1"))
MIN_SIM_THRESHOLD = float(os.getenv("MIN_SIM_THRESHOLD", "0.30"))  # used in post-filter

# Be strict about missing citations
STRICT_MODE = os.getenv("STRICT_MODE", "true").lower() in ("1", "true", "yes")
EXTRACTIVE_ONLY = os.getenv("EXTRACTIVE_ONLY", "false").lower() in ("1", "true", "yes")

# Indexing speed controls
EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "64"))
PARSE_WORKERS = int(os.getenv("PARSE_WORKERS", "1"))  # >1 enables parallel file parsing
MIN_WORDS_TO_INDEX = int(os.getenv("MIN_WORDS_TO_INDEX", "20"))

# FAISS threads (optional)
FAISS_THREADS = os.getenv("FAISS_THREADS")
try:
    if FAISS_THREADS:
        faiss.omp_set_num_threads(int(FAISS_THREADS))
except Exception:
    pass

# Embedding model (multilingual; handles Greek well)
EMBEDDING_DEVICE = os.getenv("EMBEDDING_DEVICE", "").strip().lower()  # "", "cpu", "cuda"
if not EMBEDDING_DEVICE:
    if torch is not None and getattr(torch, "cuda", None) is not None and torch.cuda.is_available():
        EMBEDDING_DEVICE = "cuda"
    else:
        EMBEDDING_DEVICE = "cpu"

embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2", device=EMBEDDING_DEVICE)

# Greek output by default
PROMPT_LANG = (
    "Απάντησε πάντα στα Ελληνικά, με φυσική και κατανοητή γλώσσα. "
    "Αν η ερώτηση είναι σε άλλη γλώσσα, μετέφρασέ την συνοπτικά και απάντησε στα Ελληνικά."
)

# Cache folder for incremental indexing
INDEX_CACHE_DIR = os.getenv("INDEX_CACHE_DIR", os.path.join(SCRIPT_DIR, ".index_cache"))
os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

# domain_data = { domain: { "chunks": [...], "index": faiss.IndexFlatIP, "meta": [...] } }
domain_data = {}

# -----------------------------
# Helpers: chunking and loading
# -----------------------------
def sliding_word_chunks(words, size=CHUNK_TOKENS, overlap=CHUNK_OVERLAP):
    """Return overlapping word chunks to improve recall."""
    if size <= 0:
        return []
    step = max(1, size - overlap)
    out = []
    for i in range(0, len(words), step):
        piece = words[i: i + size]
        if not piece:
            break
        out.append(" ".join(piece))
        if i + size >= len(words):
            break
    return out


def pdf_to_chunks_with_meta(pdf_path, chunk_size=CHUNK_TOKENS, overlap=CHUNK_OVERLAP):
    """
    Returns: (chunks, metas)
    metas: list of dicts: {"file": <name>, "page_start": int, "page_end": int, "snippet": str}
    Robust mapping: builds a word_to_page list so chunks always map to correct pages,
    even when some PDF pages extract as empty.
    """
    chunks, metas = [], []
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        per_page_text = []
        for pageno in range(len(reader.pages)):
            t = reader.pages[pageno].extract_text() or ""
            per_page_text.append(" ".join(t.split()))  # light normalization

        # Flat words array + word->page lookup
        words = []
        word_to_page = []  # same length as words; stores 0-based page number
        for pageno, text in enumerate(per_page_text):
            if text:
                w = text.split()
                words.extend(w)
                word_to_page.extend([pageno] * len(w))

        if not words:
            # If the PDF has no extractable text at all, skip it
            return [], []

        step = max(1, chunk_size - overlap)
        i = 0
        basename = os.path.basename(pdf_path)

        while i < len(words):
            j = min(i + chunk_size, len(words))
            piece_words = words[i:j]
            if not piece_words:
                break

            pages_in_slice = word_to_page[i:j]
            page_start0 = min(pages_in_slice) if pages_in_slice else 0
            page_end0 = max(pages_in_slice) if pages_in_slice else 0

            text_chunk = " ".join(piece_words)
            chunks.append(text_chunk)
            metas.append({
                "file": basename,
                "page_start": int(page_start0 + 1),  # 1-based for UI
                "page_end": int(page_end0 + 1),
                "snippet": (text_chunk[:240] + "…") if len(text_chunk) > 240 else text_chunk
            })

            if j >= len(words):
                break
            i += step

    except Exception as e:
        print(f"❌ Could not read {pdf_path}: {e}")
    return chunks, metas


def docx_to_text(docx_path):
    """Extract plain text from a .docx file."""
    doc = Document(docx_path)
    parts = []

    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)

    return "\n".join(parts)


def pptx_to_text(pptx_path):
    """Extract plain text from a .pptx file."""
    pres = Presentation(pptx_path)
    parts = []

    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text
                if text:
                    parts.append(text)

            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            parts.append(text)

    return "\n".join(parts)


# -----------------------------
# Cache utilities
# -----------------------------

def _require_cache_admin():
    """
    Cache-admin guard.

    Requires:
      - CACHE_ADMIN_KEY env var to be set (non-empty)
      - header "x-cache-admin-key" to match CACHE_ADMIN_KEY
    """
    if not CACHE_ADMIN_KEY:
        abort(403, description="Cache admin is disabled (set CACHE_ADMIN_KEY).")
    key = request.headers.get("x-cache-admin-key", "")
    if key != CACHE_ADMIN_KEY:
        abort(401, description="Unauthorized")

def _file_signature(path: str):
    """Return a stable signature used to detect changes."""
    try:
        st = os.stat(path)
        return {"size": int(st.st_size), "mtime_ns": int(getattr(st, "st_mtime_ns", int(st.st_mtime * 1e9)))}
    except Exception:
        return {"size": None, "mtime_ns": None}


def _cache_key(path: str, chunk_size: int, overlap: int):
    """
    Cache key includes file content fingerprint + chunk params.
    We avoid hashing full content for very large files by hashing (size, mtime_ns).
    If you prefer strict content hashing, set INDEX_STRICT_CONTENT_HASH=true.
    """
    strict = os.getenv("INDEX_STRICT_CONTENT_HASH", "false").lower() in ("1", "true", "yes")
    sig = _file_signature(path)
    base = f"{os.path.basename(path)}|{sig.get('size')}|{sig.get('mtime_ns')}|{chunk_size}|{overlap}|v2"
    if not strict:
        return hashlib.sha256(base.encode("utf-8")).hexdigest()

    # strict content hash (slower, but safest)
    h = hashlib.sha256()
    h.update(base.encode("utf-8"))
    try:
        with open(path, "rb") as f:
            for block in iter(lambda: f.read(1024 * 1024), b""):
                h.update(block)
    except Exception:
        pass
    return h.hexdigest()


def _domain_cache_dir(domain: str):
    d = os.path.join(INDEX_CACHE_DIR, domain)
    os.makedirs(d, exist_ok=True)
    return d


def clear_cache(domain: str | None = None) -> dict:
    """
    Delete cached embeddings/chunks.

    If domain is None: clears all domains under INDEX_CACHE_DIR.
    Returns summary dict with counts.
    """
    removed_files = 0
    removed_bytes = 0
    removed_domains = 0

    targets = []
    if domain:
        targets = [os.path.join(INDEX_CACHE_DIR, domain)]
    else:
        targets = [os.path.join(INDEX_CACHE_DIR, d) for d in os.listdir(INDEX_CACHE_DIR)]

    for t in targets:
        if not os.path.isdir(t):
            continue
        removed_domains += 1
        for root, _, files in os.walk(t):
            for fn in files:
                fp = os.path.join(root, fn)
                try:
                    st = os.stat(fp)
                    removed_bytes += int(st.st_size)
                    os.remove(fp)
                    removed_files += 1
                except Exception:
                    pass
        # remove empty dirs
        try:
            shutil.rmtree(t, ignore_errors=True)
        except Exception:
            pass

    # recreate base cache dir if needed
    os.makedirs(INDEX_CACHE_DIR, exist_ok=True)
    return {
        "domain": domain or "*",
        "removed_domains": removed_domains,
        "removed_files": removed_files,
        "removed_bytes": removed_bytes,
    }


def reload_domain(domain: str):
    """Rebuild a single domain index from BASE_FOLDER (using cache if present)."""
    folder = os.path.join(BASE_FOLDER, domain)
    if not os.path.isdir(folder):
        abort(404, description="Domain folder not found")
    build_domain_index(domain, folder)


def _cache_paths(domain: str, key: str):
    d = _domain_cache_dir(domain)
    npz_path = os.path.join(d, f"{key}.npz")
    meta_path = os.path.join(d, f"{key}.json")
    return npz_path, meta_path


def load_cached_file(domain: str, path: str, chunk_size: int, overlap: int):
    """
    If cache hit: returns (chunks, metas, embeddings) else (None, None, None).
    """
    key = _cache_key(path, chunk_size, overlap)
    npz_path, meta_path = _cache_paths(domain, key)
    if not (os.path.isfile(npz_path) and os.path.isfile(meta_path)):
        return None, None, None

    try:
        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)
        sig_now = _file_signature(path)
        if meta.get("file") != os.path.basename(path):
            return None, None, None
        if meta.get("sig") != sig_now:
            return None, None, None

        data = np.load(npz_path, allow_pickle=True)
        embeddings = data["embeddings"]
        chunks = data["chunks"].tolist()
        metas = data["metas"].tolist()
        return chunks, metas, embeddings
    except Exception as e:
        print(f"⚠️ Cache load failed for {path}: {e}")
        return None, None, None


def save_cached_file(domain: str, path: str, chunk_size: int, overlap: int, chunks, metas, embeddings):
    key = _cache_key(path, chunk_size, overlap)
    npz_path, meta_path = _cache_paths(domain, key)
    try:
        meta = {
            "file": os.path.basename(path),
            "sig": _file_signature(path),
            "chunk_size": int(chunk_size),
            "overlap": int(overlap),
            "created_utc": datetime.now(timezone.utc).isoformat(),
            "embedding_model": str(getattr(embedding_model, "model_name_or_path", "paraphrase-multilingual-MiniLM-L12-v2")),
            "device": EMBEDDING_DEVICE,
        }
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, default=str)

        # store chunks/metas as object arrays
        np.savez_compressed(
            npz_path,
            embeddings=np.asarray(embeddings, dtype=np.float32),
            chunks=np.asarray(chunks, dtype=object),
            metas=np.asarray(metas, dtype=object),
        )
    except Exception as e:
        print(f"⚠️ Cache save failed for {path}: {e}")


# -----------------------------
# File -> chunks/metas (no embeddings)
# -----------------------------
def parse_file_to_chunks(domain: str, folder: str, fname: str):
    """
    Returns:
      fname, path, chunks, metas
    """
    path = os.path.join(folder, fname)
    basename = os.path.basename(path)

    # --- PDF ---
    if fname.lower().endswith(".pdf"):
        chunks, metas = pdf_to_chunks_with_meta(path, CHUNK_TOKENS, CHUNK_OVERLAP)
        return fname, path, chunks, metas

    # --- TXT ---
    if fname.lower().endswith(".txt"):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as e:
            print(f"❌ Could not read TXT {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    # --- DOCX ---
    if fname.lower().endswith(".docx"):
        try:
            text = docx_to_text(path)
        except Exception as e:
            print(f"❌ Could not read DOCX {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    # --- PPTX ---
    if fname.lower().endswith(".pptx"):
        try:
            text = pptx_to_text(path)
        except Exception as e:
            print(f"❌ Could not read PPTX {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    return fname, path, [], []


# -----------------------------
# Domain indexing (incremental + cached)
# -----------------------------
def build_domain_index(domain, folder):
    """
    Builds (or refreshes) a domain index using:
      - Per-file caching (chunks+metas+embeddings)
      - Incremental embedding for new/changed files
      - Batched embedding encode for speed
      - Optional parallel parsing for I/O bound workloads
    """
    t0 = time.time()
    all_chunks = []
    all_metas = []
    all_embeddings = []

    files = [f for f in os.listdir(folder) if f.lower().endswith((".pdf", ".txt", ".docx", ".pptx"))]
    files.sort()

    if not files:
        print(f"⚠️ No usable data for '{domain}' in {folder}")
        return

    # 1) Try cache hits first; parse + embed only misses
    misses = []

    for fname in files:
        path = os.path.join(folder, fname)
        cached_chunks, cached_metas, cached_emb = load_cached_file(domain, path, CHUNK_TOKENS, CHUNK_OVERLAP)
        if cached_chunks is not None and cached_metas is not None and cached_emb is not None:
            all_chunks.extend(cached_chunks)
            all_metas.extend(cached_metas)
            all_embeddings.append(np.asarray(cached_emb, dtype=np.float32))
        else:
            misses.append(fname)

    # 2) Parse misses (optionally in parallel)
    parsed = []
    if misses:
        print(f"📖 Indexing domain '{domain}': {len(files)} files, {len(misses)} changed/new (cache miss).")
        if PARSE_WORKERS and PARSE_WORKERS > 1 and len(misses) > 1:
            with ThreadPoolExecutor(max_workers=PARSE_WORKERS) as ex:
                futs = {ex.submit(parse_file_to_chunks, domain, folder, fname): fname for fname in misses}
                for fut in as_completed(futs):
                    fname, path, chunks, metas = fut.result()
                    if chunks:
                        parsed.append((fname, path, chunks, metas))
        else:
            for fname in misses:
                _, path, chunks, metas = parse_file_to_chunks(domain, folder, fname)
                if chunks:
                    parsed.append((fname, path, chunks, metas))

    # Keep parsed in deterministic order
    parsed.sort(key=lambda x: x[0])

    # 3) Embed parsed misses in one (or few) batches and cache per file
    if parsed:
        # Flatten for embedding (but keep file boundaries)
        flat_chunks = []
        boundaries = []  # list of (path, start, end, metas)
        cursor = 0
        for fname, path, chunks, metas in parsed:
            start = cursor
            flat_chunks.extend(chunks)
            cursor += len(chunks)
            boundaries.append((path, start, cursor, metas))

        # Embed
        try:
            emb = embedding_model.encode(
                flat_chunks,
                batch_size=EMBED_BATCH_SIZE,
                convert_to_numpy=True,
                normalize_embeddings=True,
                show_progress_bar=(os.getenv("EMBED_PROGRESS", "false").lower() in ("1", "true", "yes"))
            )
        except TypeError:
            # Older sentence-transformers: show_progress_bar may not exist
            emb = embedding_model.encode(
                flat_chunks,
                batch_size=EMBED_BATCH_SIZE,
                convert_to_numpy=True,
                normalize_embeddings=True
            )

        emb = np.asarray(emb, dtype=np.float32)

        # Split and cache per file
        for (path, start, end, metas) in boundaries:
            file_chunks = flat_chunks[start:end]
            file_emb = emb[start:end]
            # Cache writes
            save_cached_file(domain, path, CHUNK_TOKENS, CHUNK_OVERLAP, file_chunks, metas, file_emb)

            all_chunks.extend(file_chunks)
            all_metas.extend(metas)
            all_embeddings.append(file_emb)

    # 4) Build FAISS index
    if not all_chunks or not all_embeddings:
        print(f"⚠️ No usable chunks for '{domain}' after parsing/caching.")
        return

    embeddings = np.vstack(all_embeddings).astype(np.float32)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)  # cosine on normalized vectors via dot product
    index.add(embeddings)

    domain_data[domain] = {"chunks": all_chunks, "index": index, "meta": all_metas}
    dt = time.time() - t0
    print(f"✅ Indexed {len(all_chunks)} chunks for '{domain}' on {EMBEDDING_DEVICE} in {dt:.2f}s (batch={EMBED_BATCH_SIZE}).")


# -----------------------------
# LLM call
# -----------------------------
def query_model(prompt: str) -> str:
    """Call a vLLM OpenAI-compatible Chat Completions endpoint."""
    full_prompt = f"{PROMPT_LANG}\n\n{prompt}"

    payload = {
        "model": VLLM_MODEL,
        "messages": [
            {"role": "system", "content": "You must answer strictly based on the provided QUOTES and context."},
            {"role": "user", "content": full_prompt},
        ],
        "temperature": TEMPERATURE,
        "top_p": TOP_P,
        "max_tokens": MAX_TOKENS,
        "stream": False,
    }

    resp = requests.post(VLLM_URL, json=payload, timeout=VLLM_TIMEOUT)
    resp.raise_for_status()
    data = resp.json()
    return (data.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()

def generate_alternate_queries(user_q: str, n: int = NUM_ALTERNATE_QUERIES):
    """Generate paraphrases to improve recall without broadening the topic or changing language."""
    if n <= 0:
        return [user_q]

    prompt = f"""
Παράφρασε την ερώτηση του χρήστη σε έως {n} εναλλακτικές διατυπώσεις για ΑΝΑΖΗΤΗΣΗ.
Απαιτήσεις:
- Διατήρησε το ίδιο νόημα, μην προσθέτεις/αφαιρείς υποθέματα.
- Χρησιμοποίησε την ΙΔΙΑ γλώσσα με την αρχική ερώτηση.
- Απόφυγε υπερ-γενικεύσεις και αλλαγές πεδίου.
- Κάθε παραλλαγή σε ξεχωριστή γραμμή, χωρίς αρίθμηση/περιττό κείμενο.

Ερώτηση: "{user_q}"
"""
    try:
        text = query_model(prompt)
    except Exception as e:
        print("⚠️ Alternate query generation failed:", e)
        return [user_q]

    alts = [line.strip("- ").strip() for line in text.splitlines() if line.strip()]
    all_qs = [user_q] + alts[:n]

    seen = set()
    uniq = []
    for q in all_qs:
        key = q.lower()
        if key not in seen:
            seen.add(key)
            uniq.append(q)
    return uniq


# --- citations safety net ---
def ensure_citations_from_quotes(final_answer: str, quotes: str, allowed_labels):
    """If final_answer has no [n], append the unique labels seen in QUOTES, filtered to allowed."""
    if re.search(r"\[\d+\]", final_answer):
        return final_answer

    quoted_labels = [int(n) for n in re.findall(r"\[(\d+)\]", quotes)]
    if not quoted_labels:
        return final_answer

    allowed_set = set(int(x) for x in allowed_labels)
    seen = []
    for n in quoted_labels:
        if n in allowed_set and n not in seen:
            seen.append(n)

    if not seen:
        return final_answer

    tail = "".join(f"[{n}]" for n in seen)
    return (final_answer.rstrip() + " " + tail).strip()


# --- Text normalization helpers (for robust quote validation) ---
def _norm(s: str) -> str:
    """Normalize text for resilient matching: casefold, whitespace collapse, unify quotes, unicode normalize."""
    if not s:
        return ""
    s = s.replace("“", '"').replace("”", '"').replace("’", "'").replace("…", "...")
    s = unicodedata.normalize("NFKC", s)
    s = s.lower()
    s = re.sub(r"\s+", " ", s).strip()
    return s


def word_overlap_ratio(a: str, b: str) -> float:
    """Jaccard-like overlap ratio of unique words of a found in b."""
    aw = set(_norm(a).split())
    bw = set(_norm(b).split())
    if not aw:
        return 0.0
    return len(aw & bw) / len(aw)


def auto_extract_quotes(question: str, retrieved_chunks: list, max_quotes: int = 4):
    """Extract short verbatim quotes directly from retrieved text (no model)."""
    qn = _norm(question)
    q_words = set(qn.split())

    candidates = []
    for i, ch in enumerate(retrieved_chunks, start=1):
        ch_raw = ch if isinstance(ch, str) else str(ch)
        parts = re.split(r'(?<=[\.!\?;:])\s+|\n+', ch_raw)
        for part in parts:
            p = part.strip()
            if not p:
                continue
            p = re.sub(r'\s+', ' ', p).strip()
            wc = len(p.split())
            if wc < 8 or wc > 22:
                continue
            pn = _norm(p)
            if not pn:
                continue
            pw = set(pn.split())
            if not pw:
                continue
            score = (len(pw & q_words) / max(1, len(q_words)))
            candidates.append((score, wc, p, i))

    candidates.sort(key=lambda x: (x[0], -abs(x[1] - 14)), reverse=True)

    out = []
    used = set()
    for score, wc, p, lab in candidates:
        if len(out) >= max_quotes:
            break
        key = (_norm(p), lab)
        if key in used:
            continue
        used.add(key)
        out.append((p, lab))
    return out


def messages_to_context_text(messages, max_messages=20, max_chars=6000):
    """Turn OpenAI-style chat messages into plain-text context."""
    if not isinstance(messages, list):
        return ""

    cleaned = []
    tail = messages[-max_messages:] if max_messages else messages

    for m in tail:
        if not isinstance(m, dict):
            continue
        role = (m.get("role") or "").strip().lower()
        content = (m.get("content") or "").strip()
        if not content or role not in ("system", "user", "assistant"):
            continue

        tag = {"system": "SYSTEM", "user": "USER", "assistant": "ASSISTANT"}[role]
        cleaned.append(f"{tag}: {content}")

    text = "\n".join(cleaned).strip()
    if max_chars and len(text) > max_chars:
        text = text[-max_chars:]
    return text


# -----------------------------
# Retrieval + Answer
# -----------------------------
def answer_question(domain, question, top_k=R_TOP_K, chat_messages=None):
    """
    Returns:
      final_answer: str
      sources_out: list[{"label": int, "file": str|None, "page_start": int|None,
                         "page_end": int|None, "snippet": str}]
    """
    if domain not in domain_data:
        return f"Domain '{domain}' not available.", []

    data = domain_data[domain]
    index = data["index"]
    chunks = data["chunks"]
    meta = data.get("meta")

    chat_context = messages_to_context_text(chat_messages or [])
    strict_top_k = max(int(top_k), int(os.getenv('STRICT_TOP_K', '8'))) if STRICT_MODE else int(top_k)

    if STRICT_MODE:
        queries = [question]
    else:
        queries = generate_alternate_queries(question, NUM_ALTERNATE_QUERIES)

    scored = []
    for q in queries:
        q_emb = embedding_model.encode([q], convert_to_numpy=True, normalize_embeddings=True)
        S, I = index.search(q_emb, strict_top_k)
        for sim, idx in zip(S[0].tolist(), I[0].tolist()):
            if 0 <= idx < len(chunks):
                scored.append((float(sim), int(idx)))

    if not scored:
        system_prompt = f"""
Είσαι βοηθός με ειδίκευση στο πεδίο «{domain}».
Δώσε μια σύντομη, υψηλού επιπέδου εξήγηση για το παρακάτω θέμα.
Αν δεν σχετίζεται με το πεδίο ή δεν υπάρχουν στοιχεία, απάντησε: "Δεν ξέρω."
"""
        prompt = f"{system_prompt}\n\nΕΡΩΤΗΣΗ: {question}\nΑΠΑΝΤΗΣΗ:"
        return query_model(prompt), []

    best_by_idx = {}
    for sim, idx in scored:
        best_by_idx[idx] = max(sim, best_by_idx.get(idx, -1.0))

    ranked = sorted(best_by_idx.items(), key=lambda x: x[1], reverse=True)

    if MIN_SIM_THRESHOLD is not None:
        ranked = [(idx, sim) for idx, sim in ranked if sim >= MIN_SIM_THRESHOLD]

    ranked = ranked[:strict_top_k]
    selected_idxs = [idx for idx, _ in ranked]

    retrieved = [chunks[i] for i in selected_idxs]
    allowed_labels = [str(i + 1) for i in range(len(retrieved))]
    labels_hint = ", ".join(f"[{x}]" for x in allowed_labels)

    sources_out = []
    for n, idx in enumerate(selected_idxs, start=1):
        if meta and 0 <= idx < len(meta):
            m = meta[idx]
            file = m.get("file")
            p_start = m.get("page_start")
            p_end = m.get("page_end")
            snippet = m.get("snippet") or (retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1])
        else:
            file = None
            p_start = None
            p_end = None
            snippet = retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1]

        sources_out.append({"label": n, "file": file, "page_start": p_start, "page_end": p_end, "snippet": snippet})

    if not sources_out:
        ranked = sorted(best_by_idx.items(), key=lambda x: x[1], reverse=True)[:top_k]
        selected_idxs = [idx for idx, _ in ranked]
        retrieved = [chunks[i] for i in selected_idxs]
        allowed_labels = [str(i + 1) for i in range(len(retrieved))]
        labels_hint = ", ".join(f"[{x}]" for x in allowed_labels)
        sources_out = []
        for n, idx in enumerate(selected_idxs, start=1):
            if meta and 0 <= idx < len(meta):
                m = meta[idx]
                file = m.get("file")
                p_start = m.get("page_start")
                p_end = m.get("page_end")
                snippet = m.get("snippet") or (retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1])
            else:
                file = None
                p_start = None
                p_end = None
                snippet = retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1]
            sources_out.append({"label": n, "file": file, "page_start": p_start, "page_end": p_end, "snippet": snippet})

    # -------- Pass 1: QUOTES (extractive, deterministic) --------
    extracted = auto_extract_quotes(question, retrieved, max_quotes=4)
    if not extracted:
        return ("Δεν ξέρω.", []) if STRICT_MODE else ("Δεν ξέρω.", [])

    quotes = "\n".join([f'- "{q}" [{lab}]' for (q, lab) in extracted])

    # -------- Pass 2: FINAL ANSWER --------
    if EXTRACTIVE_ONLY:
        lines = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
        final_answer = "Συνοπτικά:\n" + "\n".join(lines[:4]) if lines else "Δεν ξέρω."
    else:
        answer_system = f"""
Συντάκτης απαντήσεων, αυστηρά βασισμένων στα αποσπάσματα.
Σύνθεση ΜΟΝΟ από τα ΠΑΡΑΠΑΝΩ QUOTES, χωρίς νέα γνώση.

Κανόνες (αυστηρά):
- Απάντησε ΜΟΝΟ στην τρέχουσα ΕΡΩΤΗΣΗ. Αγνόησε παλιότερα θέματα.
- Η ΣΥΝΟΜΙΛΙΑ χρησιμοποιείται ΜΟΝΟ για συμφραζόμενα (coreference), όχι ως πηγή facts.
- Μήκος: 1–2 προτάσεις το πολύ. Χωρίς εισαγωγές/περιττά.
- Κάθε πρόταση/κουκκίδα με ισχυρισμό τελειώνει με έγκυρη ετικέτα από: {labels_hint}.
- Μην εικάζεις. Αν τα QUOTES δεν αρκούν: "Δεν ξέρω.".
"""
        final_prompt = f"""{answer_system}

ΣΥΝΟΜΙΛΙΑ (για συμφραζόμενα):
{chat_context if chat_context else "(καμία προηγούμενη συνομιλία)"}

ΕΡΩΤΗΣΗ:
{question}

QUOTES:
{quotes}

ΤΕΛΙΚΗ ΑΠΑΝΤΗΣΗ:"""
try:
        final_answer = query_model(final_prompt).strip()
except Exception as e:
        print("⚠️ vLLM call failed; returning extractive answer:", e)
        lines_ = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
        final_answer = "Συνοπτικά:\n" + "\n".join(lines_[:4]) if lines_ else "Δεν ξέρω."


    final_answer = ensure_citations_from_quotes(final_answer, quotes, allowed_labels)

    if STRICT_MODE and not re.search(r"\[\d+\]", final_answer):
        final_answer = "Δεν ξέρω."

    used_labels = sorted(set(int(n) for n in re.findall(r"\[(\d+)\]", final_answer)))
    if used_labels:
        filtered_sources = [s for s in sources_out if s["label"] in used_labels]
        order = {lab: i for i, lab in enumerate(used_labels)}
        filtered_sources.sort(key=lambda s: order.get(s["label"], 9999))
    else:
        filtered_sources = []

    return final_answer, filtered_sources


# -----------------------------
# File serving (PDFs only)
# -----------------------------
@app.route("/files/<domain>/<path:filename>", methods=["GET"])
def serve_file(domain, filename):
    if not filename.lower().endswith(".pdf"):
        abort(403)
    domain_path = os.path.join(BASE_FOLDER, domain)
    if not os.path.isdir(domain_path):
        abort(404)
    safe_path = safe_join(domain_path, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)
    return send_from_directory(domain_path, filename, as_attachment=False)


# -----------------------------
# Flask endpoints
# -----------------------------
def require_api_key():
    key = request.headers.get("x-api-key")
    if not API_KEY or API_KEY == "changeme":
        print("⚠️ WARNING: API_KEY is default. Set API_KEY env var before public exposure!")
    elif key != API_KEY:
        abort(401, description="Unauthorized")


@app.route("/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "model": MODEL_NAME,
        "domains": list(domain_data.keys()),
        "embedding_device": EMBEDDING_DEVICE,
        "chunk_tokens": CHUNK_TOKENS,
        "chunk_overlap": CHUNK_OVERLAP,
        "embed_batch_size": EMBED_BATCH_SIZE,
        "parse_workers": PARSE_WORKERS,
        "cache_dir": INDEX_CACHE_DIR,
        "cache_admin_enabled": bool(CACHE_ADMIN_KEY)
    })


@app.route("/ask", methods=["POST"])
def ask():
    require_api_key()
    data = request.get_json(force=True)
    question = (data.get("question") or "").strip()
    domain_in = data.get("domain", "default")
    chat_messages = data.get("messages") or []
    matched_domain = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), domain_in)
    print(f"Question for domain [{matched_domain}]: {question}")
    answer, sources = answer_question(matched_domain, question, chat_messages=chat_messages)
    return jsonify({"answer": answer, "sources": sources})


# -----------------------------
# DOWNLOAD PDF PAGES (excerpt)
# -----------------------------
@app.route("/files/clip", methods=["GET"])
def serve_pdf_clip():
    domain = request.args.get("domain")
    filename = request.args.get("file")
    page_from = request.args.get("from", type=int)
    page_to = request.args.get("to", type=int)

    if not (domain and filename and page_from):
        abort(400, description="Missing domain/file/from")

    if not filename.lower().endswith(".pdf"):
        abort(403)

    domain_path = os.path.join(BASE_FOLDER, domain)
    if not os.path.isdir(domain_path):
        abort(404)

    safe_path = safe_join(domain_path, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)

    if page_to is None:
        page_to = page_from

    try:
        reader = PyPDF2.PdfReader(safe_path)
        total = len(reader.pages)
        start = max(1, page_from)
        end = max(start, page_to)
        start0 = min(total, start) - 1
        end0 = min(total, end) - 1

        writer = PyPDF2.PdfWriter()
        for p in range(start0, end0 + 1):
            writer.add_page(reader.pages[p])

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
    except Exception as e:
        print("Clip error:", e)
        abort(500, description="Failed to clip PDF")

    out_name = f"{os.path.splitext(filename)[0]}_p{start}-p{end}.pdf"
    return send_file(buf, mimetype="application/pdf", as_attachment=True, download_name=out_name)




# -----------------------------
# Cache admin endpoints
# -----------------------------
@app.route("/cache/clear", methods=["POST"])
def cache_clear():
    """
    Clear cached embeddings/chunks.

    Body JSON:
      { "domain": "name" }  -> clears only that domain
      { "domain": null }    -> clears all cache
      { "reload": true }    -> rebuilds index(es) after clearing (optional)

    Auth:
      header: x-cache-admin-key must equal CACHE_ADMIN_KEY (and CACHE_ADMIN_KEY must be set)
    """
    _require_cache_admin()
    data = request.get_json(force=True, silent=True) or {}
    domain = data.get("domain")
    reload_after = bool(data.get("reload", False))

    summary = clear_cache(domain if isinstance(domain, str) and domain.strip() else None)

    if reload_after:
        if domain and isinstance(domain, str) and domain.strip():
            # also drop in-memory domain first
            domain_data.pop(domain, None)
            reload_domain(domain)
            summary["reloaded"] = [domain]
        else:
            domain_data.clear()
            load_domains()
            summary["reloaded"] = list(domain_data.keys())
    else:
        summary["reloaded"] = []

    return jsonify(summary)


@app.route("/cache/reload", methods=["POST"])
def cache_reload():
    """
    Rebuild index(es) without clearing cache.

    Body JSON:
      { "domain": "name" }  -> rebuild only that domain
      { "domain": null }    -> rebuild all domains

    Auth:
      header: x-cache-admin-key must equal CACHE_ADMIN_KEY (and CACHE_ADMIN_KEY must be set)
    """
    _require_cache_admin()
    data = request.get_json(force=True, silent=True) or {}
    domain = data.get("domain")

    if domain and isinstance(domain, str) and domain.strip():
        domain_data.pop(domain, None)
        reload_domain(domain)
        return jsonify({"reloaded": [domain]})

    domain_data.clear()
    load_domains()
    return jsonify({"reloaded": list(domain_data.keys())})


def load_domains():
    if not os.path.isdir(BASE_FOLDER):
        print(f"⚠️ BASE_FOLDER not found: {BASE_FOLDER}")
        return
    found_any = False
    for domain in os.listdir(BASE_FOLDER):
        folder = os.path.join(BASE_FOLDER, domain)
        if os.path.isdir(folder):
            build_domain_index(domain, folder)
            found_any = True
    if not found_any:
        print(f"⚠️ No domain subfolders found in {BASE_FOLDER}. Create e.g.:")
        print(os.path.join(BASE_FOLDER, "cybersecurity"))
        print(os.path.join(BASE_FOLDER, "digitalMarketing"))


load_domains()

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
