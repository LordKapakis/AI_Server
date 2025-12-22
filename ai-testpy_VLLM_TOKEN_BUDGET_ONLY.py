# ai-testpy_VLLM.py
# Stable server with:
# - legacy flat domains supported
# - nested Courses/<Course>/YEAR_X/SEM_Y/<Lesson> supported
# - year-scope escalation via scope="year"
# - ✅ accepts year alias: Courses__<Course>__YEAR_X
# - safe file serving for nested domain folders

from flask import Flask, request, jsonify, abort, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import safe_join
import os, io, re, time, json, hashlib, shutil, unicodedata
import numpy as np
import requests
import PyPDF2

try:
    import faiss
except Exception:
    faiss = None

try:
    import torch
except Exception:
    torch = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


app = Flask(__name__)
CORS(app)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# === Core Config (keep compatible) ===
DEFAULT_BASE = os.path.join(SCRIPT_DIR, "domains")
BASE_FOLDER = os.getenv("BASE_FOLDER", DEFAULT_BASE)
os.makedirs(BASE_FOLDER, exist_ok=True)

API_KEY = os.getenv("API_KEY", "changeme")

# ✅ Courses root folder name
COURSES_ROOT = os.getenv("COURSES_ROOT", "Courses")

# vLLM config
VLLM_URL = os.getenv("VLLM_URL", "http://localhost:8000/v1/chat/completions")
VLLM_MODEL = os.getenv("VLLM_MODEL", "meta-llama/Meta-Llama-3.1-8B-Instruct")
VLLM_TIMEOUT = int(os.getenv("VLLM_TIMEOUT", "120"))

MODEL_NAME = os.getenv("MODEL_NAME", "vLLM-chat")

TEMPERATURE = float(os.getenv("TEMPERATURE", "0.2"))
TOP_P = float(os.getenv("TOP_P", "0.9"))
MAX_TOKENS = int(os.getenv("MAX_TOKENS", "140"))

# vLLM token budgeting (prevents 400s when prompt is large)
VLLM_MAX_CONTEXT = int(os.getenv("VLLM_MAX_CONTEXT", "2048"))
VLLM_SAFETY_MARGIN = int(os.getenv("VLLM_SAFETY_MARGIN", "256"))
VLLM_MIN_COMPLETION = int(os.getenv("VLLM_MIN_COMPLETION", "96"))
VLLM_MAX_COMPLETION_CAP = int(os.getenv("VLLM_MAX_COMPLETION_CAP", str(MAX_TOKENS)))

# Retrieval tuning
CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "280"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "20"))
R_TOP_K = int(os.getenv("R_TOP_K", "4"))
NUM_ALTERNATE_QUERIES = int(os.getenv("NUM_ALTERNATE_QUERIES", "1"))
MIN_SIM_THRESHOLD = float(os.getenv("MIN_SIM_THRESHOLD", "0.30"))

STRICT_MODE = os.getenv("STRICT_MODE", "true").lower() in ("1", "true", "yes")
EXTRACTIVE_ONLY = os.getenv("EXTRACTIVE_ONLY", "false").lower() in ("1", "true", "yes")

EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "64"))
MIN_WORDS_TO_INDEX = int(os.getenv("MIN_WORDS_TO_INDEX", "20"))

# Index cache
INDEX_CACHE_DIR = os.getenv("INDEX_CACHE_DIR", os.path.join(SCRIPT_DIR, ".index_cache"))
os.makedirs(INDEX_CACHE_DIR, exist_ok=True)

CACHE_ADMIN_KEY = os.getenv("CACHE_ADMIN_KEY", "").strip()

SUPPORTED_EXTS = (".pdf", ".txt", ".docx", ".pptx")

# Embedding model
EMBEDDING_DEVICE = os.getenv("EMBEDDING_DEVICE", "").strip().lower()
if not EMBEDDING_DEVICE:
    if torch is not None and getattr(torch, "cuda", None) is not None and torch.cuda.is_available():
        EMBEDDING_DEVICE = "cuda"
    else:
        EMBEDDING_DEVICE = "cpu"

if SentenceTransformer is None:
    raise RuntimeError("sentence-transformers is required")
if faiss is None:
    raise RuntimeError("faiss is required")

embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2", device=EMBEDDING_DEVICE)

PROMPT_LANG = (
    "Απάντησε πάντα στα Ελληνικά, με φυσική και κατανοητή γλώσσα. "
    "Αν η ερώτηση είναι σε άλλη γλώσσα, μετέφρασέ την συνοπτικά και απάντησε στα Ελληνικά."
)

# =========================
# In-memory indexes
# =========================
# domain_data:
# { domainId: { "chunks": [...], "index": faiss.IndexFlatIP, "meta": [...],
#              "embeddings": np.ndarray, "folder": "...", "kind": "lesson|year",
#              "children": [...]} }
domain_data = {}

# domainId -> folder path
domain_folders = {}

# leaf lesson domain -> year aggregate domain (YEAR__Course__YEAR_X)
lesson_to_year = {}

# ✅ NEW: year alias -> year aggregate domain
# e.g. "Courses__Psychology__YEAR_1" -> "YEAR__Psychology__YEAR_1"
year_alias_to_year = {}


# =========================
# Helpers
# =========================
def require_api_key():
    key = request.headers.get("x-api-key")
    if not API_KEY or API_KEY == "changeme":
        print("⚠️ WARNING: API_KEY is default. Set API_KEY env var before public exposure!")
    elif key != API_KEY:
        abort(401, description="Unauthorized")


def sliding_word_chunks(words, chunk_size, overlap):
    out = []
    step = max(1, chunk_size - overlap)
    i = 0
    while i < len(words):
        out.append(" ".join(words[i:i + chunk_size]))
        i += step
    return out


def docx_to_text(docx_path):
    if docx is None:
        return ""
    d = docx.Document(docx_path)
    parts = []
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def pptx_to_text(pptx_path):
    if Presentation is None:
        return ""
    pres = Presentation(pptx_path)
    parts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                t = shape.text_frame.text
                if t:
                    parts.append(t)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            parts.append(t)
    return "\n".join(parts)


def parse_file_to_chunks(fname, folder):
    path = os.path.join(folder, fname)
    basename = fname.replace("\\", "/")

    # PDF
    if fname.lower().endswith(".pdf"):
        try:
            reader = PyPDF2.PdfReader(path)
        except Exception as e:
            print(f"❌ Could not read PDF {path}: {e}")
            return fname, path, [], []
        chunks, metas = [], []
        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            words = text.split()
            if len(words) < MIN_WORDS_TO_INDEX:
                continue
            page_chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
            for ch in page_chunks:
                chunks.append(ch)
                metas.append({
                    "file": basename,
                    "page_start": i + 1,
                    "page_end": i + 1,
                    "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
                })
        return fname, path, chunks, metas

    # TXT
    if fname.lower().endswith(".txt"):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as e:
            print(f"❌ Could not read TXT {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    # DOCX
    if fname.lower().endswith(".docx"):
        try:
            text = docx_to_text(path)
        except Exception as e:
            print(f"❌ Could not read DOCX {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    # PPTX
    if fname.lower().endswith(".pptx"):
        try:
            text = pptx_to_text(path)
        except Exception as e:
            print(f"❌ Could not read PPTX {path}: {e}")
            return fname, path, [], []
        words = (text or "").split()
        if len(words) < MIN_WORDS_TO_INDEX:
            return fname, path, [], []
        chunks = sliding_word_chunks(words, CHUNK_TOKENS, CHUNK_OVERLAP)
        metas = [{
            "file": basename,
            "page_start": None,
            "page_end": None,
            "snippet": (ch[:240] + "…") if len(ch) > 240 else ch
        } for ch in chunks]
        return fname, path, chunks, metas

    return fname, path, [], []


# =========================
# Cache
# =========================
def _require_cache_admin():
    if not CACHE_ADMIN_KEY:
        abort(403, description="Cache admin is disabled (set CACHE_ADMIN_KEY).")
    key = request.headers.get("x-cache-admin-key", "")
    if key != CACHE_ADMIN_KEY:
        abort(401, description="Unauthorized")


def _file_signature(path: str):
    try:
        st = os.stat(path)
        return {
            "size": int(st.st_size),
            "mtime_ns": int(getattr(st, "st_mtime_ns", int(st.st_mtime * 1e9))),
        }
    except Exception:
        return {"size": None, "mtime_ns": None}


def _cache_key(path: str, chunk_size: int, overlap: int):
    strict = os.getenv("INDEX_STRICT_CONTENT_HASH", "false").lower() in ("1", "true", "yes")
    sig = _file_signature(path)
    base = f"{os.path.basename(path)}|{sig.get('size')}|{sig.get('mtime_ns')}|{chunk_size}|{overlap}|v2"
    if not strict:
        return hashlib.sha256(base.encode("utf-8")).hexdigest()
    h = hashlib.sha256()
    h.update(base.encode("utf-8"))
    try:
        with open(path, "rb") as f:
            for block in iter(lambda: f.read(1024 * 1024), b""):
                h.update(block)
    except Exception:
        pass
    return h.hexdigest()


def _domain_cache_dir(domain: str):
    d = os.path.join(INDEX_CACHE_DIR, domain)
    os.makedirs(d, exist_ok=True)
    return d


def _cache_paths(domain: str, key: str):
    d = _domain_cache_dir(domain)
    npz_path = os.path.join(d, f"{key}.npz")
    meta_path = os.path.join(d, f"{key}.json")
    return npz_path, meta_path


def load_cached_file(domain: str, path: str, chunk_size: int, overlap: int):
    key = _cache_key(path, chunk_size, overlap)
    npz_path, meta_path = _cache_paths(domain, key)
    if not (os.path.isfile(npz_path) and os.path.isfile(meta_path)):
        return None, None, None
    try:
        arr = np.load(npz_path)
        emb = arr["embeddings"].astype(np.float32)
        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)
        return meta.get("chunks"), meta.get("metas"), emb
    except Exception:
        return None, None, None


def save_cached_file(domain: str, path: str, chunk_size: int, overlap: int, chunks, metas, embeddings):
    key = _cache_key(path, chunk_size, overlap)
    npz_path, meta_path = _cache_paths(domain, key)
    try:
        np.savez_compressed(npz_path, embeddings=np.asarray(embeddings, dtype=np.float32))
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump({"chunks": chunks, "metas": metas}, f, ensure_ascii=False)
    except Exception as e:
        print("⚠️ cache write failed:", e)


# =========================
# Domain discovery / IDs
# =========================
def _safe_id_from_relpath(rel: str) -> str:
    rel = rel.replace("\\", "/").strip("/")
    parts = [p.strip().replace(" ", "_") for p in rel.split("/") if p.strip()]
    return "__".join(parts)


def _find_courses_root_index(parts):
    for i, p in enumerate(parts):
        if p.lower() == COURSES_ROOT.lower():
            return i
    return None


def _extract_course_year_from_rel(rel: str):
    rel = rel.replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    i = _find_courses_root_index(parts)
    if i is None:
        return None, None
    course = parts[i + 1] if len(parts) > i + 1 else None
    year = parts[i + 2] if len(parts) > i + 2 else None
    if year and not year.upper().startswith("YEAR_"):
        return course, None
    return course, year


def _is_year_alias(domain_id: str) -> bool:
    # "Courses__Psychology__YEAR_1"
    if not isinstance(domain_id, str):
        return False
    parts = domain_id.split("__")
    return len(parts) == 3 and parts[0] == "Courses" and parts[2].upper().startswith("YEAR_")


def discover_domain_folders():
    found = []
    base_abs = os.path.abspath(BASE_FOLDER)
    cache_abs = os.path.abspath(INDEX_CACHE_DIR)

    for root, dirs, files in os.walk(BASE_FOLDER):
        root_abs = os.path.abspath(root)
        if root_abs.startswith(cache_abs):
            continue

        usable = [f for f in files if f.lower().endswith(SUPPORTED_EXTS)]
        if not usable:
            continue

        rel = os.path.relpath(root, BASE_FOLDER).replace("\\", "/")
        domain_id = _safe_id_from_relpath(rel)
        found.append((domain_id, root, rel))

    return found


# =========================
# Index building
# =========================
def build_domain_index(domain: str, folder: str):
    t0 = time.time()

    files = [f for f in os.listdir(folder) if f.lower().endswith(SUPPORTED_EXTS)]
    files.sort()
    if not files:
        print(f"⚠️ No usable files for '{domain}' in {folder}")
        return

    all_chunks = []
    all_metas = []
    all_embeddings_parts = []

    misses = []
    for f in files:
        path = os.path.join(folder, f)
        ch, metas, emb = load_cached_file(domain, path, CHUNK_TOKENS, CHUNK_OVERLAP)
        if ch is not None and metas is not None and emb is not None:
            all_chunks.extend(ch)
            all_metas.extend(metas)
            all_embeddings_parts.append(np.asarray(emb, dtype=np.float32))
        else:
            misses.append(f)

    parsed = []
    for f in misses:
        fname, path, chunks, metas = parse_file_to_chunks(f, folder)
        if chunks and metas:
            parsed.append((path, chunks, metas))

    if parsed:
        flat_chunks = []
        boundaries = []
        cursor = 0

        for path, chunks, metas in parsed:
            start = cursor
            flat_chunks.extend(chunks)
            cursor += len(chunks)
            boundaries.append((path, start, cursor, metas))

        try:
            emb = embedding_model.encode(
                flat_chunks,
                batch_size=EMBED_BATCH_SIZE,
                convert_to_numpy=True,
                normalize_embeddings=True,
                show_progress_bar=(os.getenv("EMBED_PROGRESS", "false").lower() in ("1", "true", "yes"))
            )
        except TypeError:
            emb = embedding_model.encode(
                flat_chunks,
                batch_size=EMBED_BATCH_SIZE,
                convert_to_numpy=True,
                normalize_embeddings=True
            )

        emb = np.asarray(emb, dtype=np.float32)

        for (path, start, end, metas) in boundaries:
            file_chunks = flat_chunks[start:end]
            file_emb = emb[start:end]
            save_cached_file(domain, path, CHUNK_TOKENS, CHUNK_OVERLAP, file_chunks, metas, file_emb)

            all_chunks.extend(file_chunks)
            all_metas.extend(metas)
            all_embeddings_parts.append(file_emb)

    if not all_chunks or not all_embeddings_parts:
        print(f"⚠️ No usable chunks for '{domain}'")
        return

    embeddings = np.vstack(all_embeddings_parts).astype(np.float32)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)

    domain_data[domain] = {
        "chunks": all_chunks,
        "meta": all_metas,
        "embeddings": embeddings,
        "index": index,
        "folder": folder,
        "kind": "lesson",
    }
    domain_folders[domain] = folder

    dt = time.time() - t0
    print(f"✅ Indexed {len(all_chunks)} chunks for '{domain}' in {dt:.2f}s.")


def build_year_aggregates():
    """
    Builds YEAR__<Course>__YEAR_X aggregates from lesson domains under Courses/<Course>/YEAR_X/...
    Also builds alias map:
      Courses__<Course>__YEAR_X  ->  YEAR__<Course>__YEAR_X
    """
    year_map = {}

    # reset maps (safe)
    year_alias_to_year.clear()

    for d, info in list(domain_data.items()):
        if info.get("kind") != "lesson":
            continue

        folder = info.get("folder")
        if not folder:
            continue

        rel = os.path.relpath(folder, BASE_FOLDER).replace("\\", "/")
        course, year = _extract_course_year_from_rel(rel)
        if not course or not year:
            continue

        course_id = course.replace(" ", "_")
        year_id = year.replace(" ", "_")

        year_domain = f"YEAR__{course_id}__{year_id}"
        year_alias = f"Courses__{course_id}__{year_id}"  # ✅ what frontend sends

        year_map.setdefault(year_domain, []).append(d)
        lesson_to_year[d] = year_domain
        year_alias_to_year[year_alias] = year_domain

    for year_domain, children in year_map.items():
        chunks_all = []
        meta_all = []
        emb_parts = []

        for child in children:
            child_info = domain_data.get(child)
            if not child_info:
                continue
            chunks_all.extend(child_info.get("chunks", []))
            meta_all.extend(child_info.get("meta", []))
            emb = child_info.get("embeddings")
            if emb is not None and len(emb) > 0:
                emb_parts.append(np.asarray(emb, dtype=np.float32))

        if not chunks_all or not emb_parts:
            continue

        embeddings = np.vstack(emb_parts).astype(np.float32)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)

        domain_data[year_domain] = {
            "chunks": chunks_all,
            "meta": meta_all,
            "embeddings": embeddings,
            "index": index,
            "folder": None,
            "kind": "year",
            "children": children,
        }

        print(f"🧩 YEAR index '{year_domain}' built from {len(children)} lesson domains.")


def load_domains():
    domain_data.clear()
    domain_folders.clear()
    lesson_to_year.clear()
    year_alias_to_year.clear()

    discovered = discover_domain_folders()

    for domain_id, folder, rel in discovered:
        build_domain_index(domain_id, folder)

    build_year_aggregates()


# =========================
# LLM
# =========================
def _approx_token_count(text: str) -> int:
    if not text:
        return 0
    return max(1, int(len(text) / 4))


def _truncate_text_to_token_budget(text: str, max_tokens: int) -> str:
    if max_tokens <= 0:
        return ""
    approx_tokens = _approx_token_count(text)
    if approx_tokens <= max_tokens:
        return text
    keep_chars = int(len(text) * (max_tokens / max(approx_tokens, 1)) * 0.98)
    keep_chars = max(0, min(len(text), keep_chars))
    if keep_chars <= 0:
        return text[-min(len(text), 2000):]
    return text[-keep_chars:]


def query_model(prompt: str) -> str:
    full_prompt = f"{PROMPT_LANG}\n\n{prompt}"

    input_tokens = _approx_token_count(full_prompt) + 64
    budget = VLLM_MAX_CONTEXT - VLLM_SAFETY_MARGIN
    remaining_for_completion = max(0, budget - input_tokens)

    max_completion = min(VLLM_MAX_COMPLETION_CAP, max(VLLM_MIN_COMPLETION, remaining_for_completion))

    if remaining_for_completion < VLLM_MIN_COMPLETION:
        max_input_allowed = max(128, budget - VLLM_MIN_COMPLETION)
        full_prompt = _truncate_text_to_token_budget(full_prompt, max_input_allowed)

    payload = {
        "model": VLLM_MODEL,
        "messages": [
            {"role": "system", "content": "You must answer strictly based on the provided QUOTES and context."},
            {"role": "user", "content": full_prompt},
        ],
        "temperature": TEMPERATURE,
        "top_p": TOP_P,
        "max_tokens": int(max_completion),
        "stream": False,
    }
    resp = requests.post(VLLM_URL, json=payload, timeout=VLLM_TIMEOUT)
    resp.raise_for_status()
    data = resp.json()
    return (data.get("choices", [{}])[0].get("message", {}).get("content") or "").strip()


def generate_alternate_queries(user_q: str, n: int = NUM_ALTERNATE_QUERIES):
    if n <= 0:
        return [user_q]
    prompt = f"""
Παράφρασε την ερώτηση του χρήστη σε έως {n} εναλλακτικές διατυπώσεις για ΑΝΑΖΗΤΗΣΗ.
- Κάθε παραλλαγή σε ξεχωριστή γραμμή, χωρίς αρίθμηση/περιττό κείμενο.
Ερώτηση: "{user_q}"
"""
    try:
        text = query_model(prompt)
    except Exception as e:
        print("⚠️ Alternate query generation failed:", e)
        return [user_q]

    alts = [line.strip("- ").strip() for line in text.splitlines() if line.strip()]
    all_qs = [user_q] + alts[:n]

    seen = set()
    uniq = []
    for q in all_qs:
        key = q.lower()
        if key not in seen:
            seen.add(key)
            uniq.append(q)
    return uniq


# =========================
# Strict citation/extractive logic
# =========================
def _norm(s: str) -> str:
    if not s:
        return ""
    s = s.replace("“", '"').replace("”", '"').replace("’", "'").replace("…", "...")
    s = unicodedata.normalize("NFKC", s)
    s = s.lower()
    s = re.sub(r"\s+", " ", s).strip()
    return s


def auto_extract_quotes(question: str, retrieved_chunks: list, max_quotes: int = 4):
    qn = _norm(question)
    q_words = set(qn.split())
    candidates = []
    for i, ch in enumerate(retrieved_chunks, start=1):
        ch_raw = ch if isinstance(ch, str) else str(ch)
        parts = re.split(r'(?<=[\.!\?;:])\s+|\n+', ch_raw)
        for part in parts:
            p = part.strip()
            if not p:
                continue
            p = re.sub(r'\s+', ' ', p).strip()
            wc = len(p.split())
            if wc < 8 or wc > 22:
                continue
            pn = _norm(p)
            if not pn:
                continue
            pw = set(pn.split())
            if not pw:
                continue
            score = (len(pw & q_words) / max(1, len(q_words)))
            candidates.append((score, wc, p, i))

    candidates.sort(key=lambda x: (x[0], -abs(x[1] - 14)), reverse=True)
    out = []
    used = set()
    for score, wc, p, lab in candidates:
        if len(out) >= max_quotes:
            break
        key = (_norm(p), lab)
        if key in used:
            continue
        used.add(key)
        out.append((p, lab))
    return out


def ensure_citations_from_quotes(final_answer: str, quotes: str, allowed_labels):
    if re.search(r"\[\d+\]", final_answer):
        return final_answer
    quoted_labels = [int(n) for n in re.findall(r"\[(\d+)\]", quotes)]
    if not quoted_labels:
        return final_answer
    allowed_set = set(int(x) for x in allowed_labels)
    seen = []
    for n in quoted_labels:
        if n in allowed_set and n not in seen:
            seen.append(n)
    if not seen:
        return final_answer
    tail = "".join(f"[{n}]" for n in seen)
    return (final_answer.rstrip() + " " + tail).strip()


def messages_to_context_text(messages, max_messages=20, max_chars=6000):
    if not isinstance(messages, list):
        return ""
    cleaned = []
    tail = messages[-max_messages:] if max_messages else messages
    for m in tail:
        if not isinstance(m, dict):
            continue
        role = (m.get("role") or "").strip().lower()
        content = (m.get("content") or "").strip()
        if not content or role not in ("system", "user", "assistant"):
            continue
        tag = {"system": "SYSTEM", "user": "USER", "assistant": "ASSISTANT"}[role]
        cleaned.append(f"{tag}: {content}")
    text = "\n".join(cleaned).strip()
    if max_chars and len(text) > max_chars:
        text = text[-max_chars:]
    return text


# =========================
# Retrieval + Answer
# =========================
def answer_question(domain, question, top_k=R_TOP_K, chat_messages=None):
    if domain not in domain_data:
        return f"Domain '{domain}' not available.", []

    data = domain_data[domain]
    index = data["index"]
    chunks = data["chunks"]
    meta = data.get("meta")

    chat_context = messages_to_context_text(chat_messages or [])
    strict_top_k = max(int(top_k), int(os.getenv('STRICT_TOP_K', '8'))) if STRICT_MODE else int(top_k)

    if STRICT_MODE:
        queries = [question]
    else:
        queries = generate_alternate_queries(question, NUM_ALTERNATE_QUERIES)

    scored = []
    for q in queries:
        q_emb = embedding_model.encode([q], convert_to_numpy=True, normalize_embeddings=True)
        S, I = index.search(q_emb, strict_top_k)
        for sim, idx in zip(S[0].tolist(), I[0].tolist()):
            if 0 <= idx < len(chunks):
                scored.append((float(sim), int(idx)))

    if not scored:
        system_prompt = f"""
Είσαι βοηθός με ειδίκευση στο πεδίο «{domain}».
Αν δεν υπάρχουν στοιχεία, απάντησε: "Δεν ξέρω."
"""
        prompt = f"{system_prompt}\n\nΕΡΩΤΗΣΗ: {question}\nΑΠΑΝΤΗΣΗ:"
        return query_model(prompt), []

    best_by_idx = {}
    for sim, idx in scored:
        best_by_idx[idx] = max(sim, best_by_idx.get(idx, -1.0))

    ranked = sorted(best_by_idx.items(), key=lambda x: x[1], reverse=True)
    if MIN_SIM_THRESHOLD is not None:
        ranked = [(idx, sim) for idx, sim in ranked if sim >= MIN_SIM_THRESHOLD]

    ranked = ranked[:strict_top_k]
    selected_idxs = [idx for idx, _ in ranked]
    retrieved = [chunks[i] for i in selected_idxs]
    allowed_labels = [str(i + 1) for i in range(len(retrieved))]
    labels_hint = ", ".join(f"[{x}]" for x in allowed_labels)

    sources_out = []
    for n, idx in enumerate(selected_idxs, start=1):
        if meta and 0 <= idx < len(meta):
            m = meta[idx]
            file = m.get("file")
            p_start = m.get("page_start")
            p_end = m.get("page_end")
            snippet = m.get("snippet") or (retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1])
        else:
            file = None
            p_start = None
            p_end = None
            snippet = retrieved[n - 1][:240] + "…" if len(retrieved[n - 1]) > 240 else retrieved[n - 1]
        sources_out.append({"label": n, "file": file, "page_start": p_start, "page_end": p_end, "snippet": snippet})

    extracted = auto_extract_quotes(question, retrieved, max_quotes=4)
    if not extracted:
        return ("Δεν ξέρω.", []) if STRICT_MODE else ("Δεν ξέρω.", [])

    quotes = "\n".join([f'- "{q}" [{lab}]' for (q, lab) in extracted])

    if EXTRACTIVE_ONLY:
        lines = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
        final_answer = "Συνοπτικά:\n" + "\n".join(lines[:4]) if lines else "Δεν ξέρω."
    else:
        answer_system = f"""
Συντάκτης απαντήσεων, αυστηρά βασισμένων στα αποσπάσματα.

Κανόνες (αυστηρά):
- Απάντησε ΜΟΝΟ στην τρέχουσα ΕΡΩΤΗΣΗ.
- Η ΣΥΝΟΜΙΛΙΑ χρησιμοποιείται ΜΟΝΟ για συμφραζόμενα (coreference), όχι ως πηγή facts.
- Μήκος: 1–2 προτάσεις το πολύ.
- Κάθε πρόταση με ισχυρισμό τελειώνει με έγκυρη ετικέτα από: {labels_hint}.
- Μην εικάζεις. Αν τα QUOTES δεν αρκούν: "Δεν ξέρω.".
"""
        final_prompt = f"""{answer_system}

ΣΥΝΟΜΙΛΙΑ:
{chat_context if chat_context else "(καμία προηγούμενη συνομιλία)"}

ΕΡΩΤΗΣΗ:
{question}

QUOTES:
{quotes}

ΤΕΛΙΚΗ ΑΠΑΝΤΗΣΗ:"""
        try:
            final_answer = query_model(final_prompt).strip()
        except Exception as e:
            print("⚠️ vLLM call failed; returning extractive answer:", e)
            lines_ = [ln.strip() for ln in quotes.splitlines() if ln.strip().startswith("- ")]
            final_answer = "Συνοπτικά:\n" + "\n".join(lines_[:4]) if lines_ else "Δεν ξέρω."

    final_answer = ensure_citations_from_quotes(final_answer, quotes, allowed_labels)
    if STRICT_MODE and not re.search(r"\[\d+\]", final_answer):
        final_answer = "Δεν ξέρω."

    used_labels = sorted(set(int(n) for n in re.findall(r"\[(\d+)\]", final_answer)))
    if used_labels:
        filtered_sources = [s for s in sources_out if s["label"] in used_labels]
        order = {lab: i for i, lab in enumerate(used_labels)}
        filtered_sources.sort(key=lambda s: order.get(s["label"], 9999))
    else:
        filtered_sources = []

    return final_answer, filtered_sources


# =========================
# Domain folder resolving
# =========================
def _resolve_domain_folder(domain: str):
    if domain in domain_folders:
        return domain_folders[domain]
    d2 = next((k for k in domain_folders.keys() if k.lower() == (domain or "").lower()), None)
    if d2:
        return domain_folders[d2]
    return None


@app.route("/files/<domain>/<path:filename>", methods=["GET"])
def serve_file(domain, filename):
    if not filename.lower().endswith(".pdf"):
        abort(403)
    folder = _resolve_domain_folder(domain)
    if not folder or not os.path.isdir(folder):
        abort(404)
    safe_path = safe_join(folder, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)
    return send_from_directory(folder, filename, as_attachment=False)


@app.route("/files/clip", methods=["GET"])
def serve_pdf_clip():
    domain = request.args.get("domain")
    filename = request.args.get("file")
    page_from = request.args.get("from", type=int)
    page_to = request.args.get("to", type=int)

    if not (domain and filename and page_from):
        abort(400, description="Missing domain/file/from")
    if not filename.lower().endswith(".pdf"):
        abort(403)

    folder = _resolve_domain_folder(domain)
    if not folder or not os.path.isdir(folder):
        abort(404)

    safe_path = safe_join(folder, filename)
    if not safe_path or not os.path.isfile(safe_path):
        abort(404)

    if page_to is None:
        page_to = page_from

    try:
        reader = PyPDF2.PdfReader(safe_path)
        total = len(reader.pages)
        start = max(1, page_from)
        end = max(start, page_to)
        start0 = min(total, start) - 1
        end0 = min(total, end) - 1

        writer = PyPDF2.PdfWriter()
        for p in range(start0, end0 + 1):
            writer.add_page(reader.pages[p])

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
    except Exception as e:
        print("Clip error:", e)
        abort(500, description="Failed to clip PDF")

    out_name = f"{os.path.splitext(os.path.basename(filename))[0]}_p{start}-p{end}.pdf"
    return send_file(buf, mimetype="application/pdf", as_attachment=True, download_name=out_name)


# =========================
# Health + Ask (scope)
# =========================
@app.route("/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "model": MODEL_NAME,
        "embedding_device": EMBEDDING_DEVICE,
        "domains_count": len(domain_data),
        "domains": list(domain_data.keys())[:200],
        "year_domains": [d for d, v in domain_data.items() if v.get("kind") == "year"],
        "year_aliases_sample": list(year_alias_to_year.items())[:50],
        "courses_root": COURSES_ROOT,
    })


@app.route("/ask", methods=["POST"])
def ask():
    require_api_key()
    data = request.get_json(force=True) or {}

    question = (data.get("question") or "").strip()
    domain_in = (data.get("domain") or "").strip()
    scope = (data.get("scope") or "lesson").strip().lower()  # lesson | year
    chat_messages = data.get("messages") or []

    if not question:
        return jsonify({"answer": "Δεν δόθηκε ερώτηση.", "sources": []})

    # match domain case-insensitively
    matched_domain = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), domain_in)

    # ✅ scope escalation (supports BOTH lesson id and year-alias id)
    if scope == "year":
        # 1) if frontend sends "Courses__Course__YEAR_X", map to year aggregate
        alias_key = domain_in.replace(" ", "_")
        if alias_key in year_alias_to_year and year_alias_to_year[alias_key] in domain_data:
            matched_domain = year_alias_to_year[alias_key]
        else:
            # 2) if it's a lesson id, map to its year aggregate
            yd = lesson_to_year.get(matched_domain)
            if yd and yd in domain_data:
                matched_domain = yd
            else:
                # 3) if frontend already sent YEAR__..., keep it
                yd2 = next((d for d in domain_data.keys() if d.lower() == domain_in.lower()), None)
                if yd2 and domain_data.get(yd2, {}).get("kind") == "year":
                    matched_domain = yd2

    print(f"Question for domain [{matched_domain}] (scope={scope}): {question}")
    answer, sources = answer_question(matched_domain, question, chat_messages=chat_messages)

    return jsonify({"answer": answer, "sources": sources, "resolvedDomain": matched_domain})


# =========================
# Cache admin (optional)
# =========================
@app.route("/cache/clear", methods=["POST"])
def cache_clear():
    _require_cache_admin()
    data = request.get_json(force=True, silent=True) or {}
    domain = data.get("domain")
    reload_after = bool(data.get("reload", False))

    summary = clear_cache(domain if isinstance(domain, str) and domain.strip() else None)

    if reload_after:
        load_domains()
        summary["reloaded"] = True
    else:
        summary["reloaded"] = False

    return jsonify(summary)


def clear_cache(domain=None):
    removed_files = 0
    removed_bytes = 0
    removed_domains = 0

    targets = []
    if domain:
        targets = [os.path.join(INDEX_CACHE_DIR, domain)]
    else:
        targets = [os.path.join(INDEX_CACHE_DIR, d) for d in os.listdir(INDEX_CACHE_DIR)]

    for t in targets:
        if not os.path.isdir(t):
            continue
        removed_domains += 1
        for root, _, files in os.walk(t):
            for fn in files:
                fp = os.path.join(root, fn)
                try:
                    st = os.stat(fp)
                    removed_bytes += int(st.st_size)
                    os.remove(fp)
                    removed_files += 1
                except Exception:
                    pass
        try:
            shutil.rmtree(t, ignore_errors=True)
        except Exception:
            pass

    os.makedirs(INDEX_CACHE_DIR, exist_ok=True)
    return {
        "domain": domain or "*",
        "removed_domains": removed_domains,
        "removed_files": removed_files,
        "removed_bytes": removed_bytes,
    }


# =========================
# Boot
# =========================
print("🔄 Loading domains...")
load_domains()
print(f"✅ Ready. Loaded {len(domain_data)} domains (including YEAR aggregates).")

if __name__ == "__main__":
    port = int(os.getenv("PORT", "5000"))
    app.run(host="0.0.0.0", port=port)
